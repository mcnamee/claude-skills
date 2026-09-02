#!/usr/bin/env python3
# -*- coding: utf-8 -*-
r"""
powerpoint.py (v3.0.0) - A single-file MCP (Model Context Protocol) stdio server
that builds PowerPoint .pptx decks, optionally from your own template, and
audits them against Guy Kawasaki's 10/20/30 rule.

It follows the same create -> build -> save workflow as word.py
(powerpoint_create ... powerpoint_save), with the .pptx engine provided by the
Python library `python-pptx`.

THE ONE IDEA THIS SERVER IS BUILT AROUND
    A deck inherits its look from its TEMPLATE - the slide master, the slide
    layouts and the theme. So this server never sets a font, a colour or a
    position: it puts text into the LAYOUT PLACEHOLDERS the template already
    defines, and lets the template style it. That is what "sticks to the
    template styles" means in practice, and it is why there is no
    powerpoint_set_font tool. If a deck looks wrong, the fix is a different
    layout, not a hard-coded typeface.

WHAT IT CAN DO
    - Create a NEW .pptx (powerpoint_create) in the presentations folder.
      With 'template' set to an existing .pptx/.potx, the new deck inherits that
      file's masters, layouts, theme, fonts and colours; the template file is
      only READ and is never modified. Without one, the stock Office template is
      used.
    - Open an EXISTING .pptx (powerpoint_open) to read it or to inspect a
      template before drafting. You can pass just the file NAME
      (e.g. "Deck Template.pptx") - paths are resolved relative to the
      presentation root, so no absolute path is needed; a bare name is even
      found in a subfolder. If no exact match exists, a FUZZY name match is
      tried (so "quarterly deck" opens "Quarterly Deck 2024.pptx") and the
      result is flagged (fuzzy_matched) so the caller can confirm the pick.
      List what is available with powerpoint_list_presentations.
    - Keep blank templates in a folder of their own with --templates-dir
      (by default C:\Eva\templates\powerpoint). It is a READ-ONLY second root:
      its files can be listed, opened and used as the base for
      powerpoint_create, but every attempt to SAVE over one is refused, so a
      template cannot be turned into someone's half-finished deck.
    - Report what the template actually offers (powerpoint_list_layouts): every
      slide layout, its placeholders (idx, type, name) and - the useful part -
      the EFFECTIVE font size each placeholder will render at, resolved through
      the full PowerPoint inheritance chain. A `recommended` block names the
      best layout in THIS template for a title slide, a section divider,
      bulleted content, two-column content and a bare title.
    - BUILD A WHOLE DECK IN ONE CALL (powerpoint_add_slides): an ordered list of
      slides, each with its layout, title, subtitle, bullets, extra placeholder
      fills, a table and speaker notes, appended in exactly the order given.
      This is the tool to build a deck with, and the only one whose slide order
      is guaranteed: see "DECK ORDER" below for why that matters.
    - Add a single slide onto a chosen layout (powerpoint_add_slide), filling the
      layout's placeholders by name or index: a title, a subtitle, and bullets
      with real outline levels. Unfilled content placeholders are removed by
      default so the deck has no "Click to add text" prompts left in it.
    - Edit what is there: powerpoint_set_placeholder, powerpoint_add_bullets,
      powerpoint_add_table, powerpoint_set_notes (speaker notes),
      powerpoint_delete_slide and powerpoint_move_slide.
    - Read a deck back (powerpoint_get_content): every slide's layout, its
      placeholder text, its tables and its speaker notes.
    - AUDIT a deck against the 10/20/30 rule (powerpoint_review): slide count vs
      10, estimated speaking time vs 20 minutes (from the speaker-note word
      count), and every run of text whose EFFECTIVE size is below 30pt - with
      the slide, the shape and the resolved size, so each finding is actionable.
      The thresholds are config constants, so a house style of, say, 24pt is a
      one-line change rather than a fork.
    - On open, create and save, optionally MIRROR the deck to Markdown into a
      knowledge-base folder for a local RAG index (the same idea as word.py's
      --kb-dir): slide titles become headings, bullets become lists, tables
      become Markdown tables and speaker notes are included. Files are named
      'PowerPoint - <name>.md' and overwritten each time.

DECK ORDER - WHY powerpoint_add_slides EXISTS
    powerpoint_add_slide appends its slide to the END of the deck. That makes
    the deck's slide order simply the order in which the CALLS REACH this
    server - and an MCP client is free to dispatch independent tool calls
    concurrently, in which case their arrival order is NOT the order the model
    wrote them. Building a ten-slide deck as ten separate add_slide calls is
    therefore a race, and it can come back with the slides shuffled. (word.py
    had the same flaw with its one-block-per-call append tools, where it showed
    up as every heading bunched together ahead of the body text.)

    The follow-up tools make it worse rather than better: powerpoint_set_notes,
    powerpoint_set_placeholder, powerpoint_add_bullets and powerpoint_add_table
    all address a slide by slide_index, so a caller that ASSUMES "the third
    slide I created is index 2" writes its notes onto whichever slide actually
    landed there.

    Nothing inside a single add_slide call can detect this - the server sees
    "append a slide" and has no idea what was meant to come before it. So the
    fix is to stop splitting the sequence: powerpoint_add_slides takes the whole
    ordered list of slides in ONE request, which cannot be reordered by
    anything, and each entry carries that slide's own extra placeholder fills,
    table and notes so no slide_index is ever guessed. Build decks with it; keep
    powerpoint_add_slide for a genuinely single slide.

    As a safety net for callers still chaining add_slide, three or more calls
    landing on the same session within APPEND_BURST_SECONDS - the signature of a
    parallel batch, not of separate conversational turns - adds an
    `order_warning` to the result, so a shuffle is reported instead of silently
    shipped.

WHY THE FONT-SIZE AUDIT IS NOT JUST run.font.size
    In a well-built deck almost NO run carries an explicit size: it is inherited
    from the layout, then the master, then the master's txStyles, then the
    presentation default. python-pptx reports `run.font.size is None` for every
    one of those, so the naive check finds nothing wrong with a deck set
    entirely in 12pt. This server resolves the size the way PowerPoint does, in
    order:
        1. the run's own a:rPr/@sz
        2. the paragraph's a:pPr/a:defRPr/@sz
        3. the shape's a:lstStyle for that outline level
        4. the matching placeholder on the SLIDE LAYOUT (same idx)
        5. the matching placeholder on the SLIDE MASTER (same type)
        6. the master's p:txStyles - titleStyle / bodyStyle / otherStyle -
           for that outline level
        7. the presentation's p:defaultTextStyle
        8. failing all of that, POWERPOINT_DEFAULT_FONT_PT (18pt, PowerPoint's
           own fallback)
    Every finding reports which of those steps supplied the number ('source'),
    so "why is this 28pt?" has an answer instead of a shrug. Worth knowing: the
    STOCK Office template already sets body level 2 to 28pt and level 3 to 24pt,
    so a default-template deck fails the 30-point rule the moment you use a
    sub-bullet. That is a real finding, not a false positive.

VERIFIED BEFORE DELIVERY
    The font-resolution chain above is the part most worth distrusting, so it
    was checked against an INDEPENDENT renderer rather than against itself. A
    deck was built from a template whose master sets titleStyle to 40pt and
    bodyStyle level 1 to 32pt, then exported to PDF by LibreOffice Impress and
    the rendered sizes read back out of the PDF:

        element                 this server      LibreOffice rendered
        title                   40pt             40.0pt
        bullet, level 1         32pt             32.0pt
        bullet, level 2         28pt             28.0pt   <- flagged, correctly
        table cells             (unmeasured)     18.0pt

    So the resolved sizes are the sizes that actually render, and the one run
    the 30-point audit flagged is genuinely the one that was too small. Also
    verified: the saved package carries the template's own master/theme parts
    and its slides contain NO hard-coded sz= attributes (the text really does
    inherit rather than being baked in), and deleting slides leaves no orphaned
    parts behind in the package.

    Table text is the honest gap. It is sized by the table style in
    tableStyles.xml, which this server does not read, so it is reported as
    unmeasured with a clearly-labelled estimate rather than counted as a
    violation. In the stock theme that estimate happened to match the render
    exactly; in a template whose table style differs from its otherStyle it
    would not, which is precisely why it is not asserted as a measurement.

    The ONE thing not verifiable off your network is Microsoft PowerPoint's own
    rendering. Open one generated deck in PowerPoint on your endpoint and
    confirm the template's fonts and colours look right before relying on this
    for anything that matters.

WHAT IT CANNOT DO
    - Charts, SmartArt, animations, transitions and slide-master EDITING. Only
      python-pptx's supported surface is used; a chart needs data plumbing this
      server deliberately does not carry.
    - Pictures. Nothing here reads image files, by design - the sandbox stays a
      text-only surface. Use a layout with a PICTURE placeholder and drop the
      image in by hand.
    - Changing a template's theme, fonts or colours. That is the template's job
      and editing it here would defeat the point.
    - .odp and Google Slides.

=============================================================================
 DEPENDENCIES
=============================================================================
    python-pptx   (import name: pptx)   -> the .pptx engine
      +- lxml            (compiled C extension, pulled in by python-pptx)
      +- Pillow          (image handling; python-pptx imports it at load)
      +- XlsxWriter      (only used for charts, still a hard dependency)
      +- typing_extensions

    This is the SAME lxml the `word` plugin needs, so if word.py already runs on
    an interpreter, most of the work is done.

    AIRGAPPED WINDOWS INSTALL (sideload wheels, same pattern as python-docx):
      1. On an internet-connected box, download wheels for your EXACT
         interpreter. lxml and Pillow are compiled, so their wheels must match
         Python version + architecture, e.g. for CPython 3.11 64-bit:
             lxml-<ver>-cp311-cp311-win_amd64.whl
             pillow-<ver>-cp311-cp311-win_amd64.whl
         python-pptx, XlsxWriter and typing_extensions are pure-Python.
         Grab the correct set in one go, ideally ON the target Python version:
             python -m pip download python-pptx -d .\wheels
      2. Transfer the .\wheels folder to the airgapped endpoint.
      3. Install with the SAME interpreter the MCP client will launch (use -m
         pip so the interpreter and pip cannot drift apart). PowerShell - the
         backtick is the line continuation, and "&" is required because the
         interpreter path is quoted:
             & "C:\path\to\python.exe" -m pip install --no-index `
                 --find-links .\wheels python-pptx
      4. Confirm the interpreter can see it:
             & "C:\path\to\python.exe" -c "import pptx; print(pptx.__version__)"

    If your corporate mirror proxies PyPI you may be able to skip the manual
    download and just run:  python -m pip install python-pptx
    A compiled-wheel/interpreter mismatch is the #1 failure here; this server
    logs sys.executable and every dependency version at startup so a mismatch
    is obvious.

=============================================================================
 VALIDATE BEFORE WIRING IN  (single-transfer sanity check)
=============================================================================
    Run the built-in self-test. It builds a temp template, creates a deck from
    it, adds slides, saves, reopens and audits - then prints PASS/FAIL. No
    arguments, no network, no side files left behind:
        & "C:\path\to\python.exe" powerpoint.py --check

    Expected tail of output on success:
        [check] 10/20/30 review: PASS
        [check] ALL CHECKS PASSED

=============================================================================
 INSTALLING INTO CLAUDE CODE
=============================================================================
    This server ships as the "powerpoint" Claude Code plugin (its manifest is
    .claude-plugin/plugin.json next to this file), so the normal install is:

        /plugin marketplace add C:\path\to\claude-skills
        /plugin install powerpoint@mcnamee-claude-skills

    Claude Code then prompts for each setting below and for the Python
    interpreter - use the SAME python.exe you installed the wheels with.
    PYTHONUTF8=1 (so Windows cp1252 cannot corrupt the stdio JSON stream) is
    set for you by the manifest.

    To register the server by hand instead:

        claude mcp add powerpoint --scope user -e PYTHONUTF8=1 -- C:\path\to\python.exe C:\path\to\powerpoint.py

    Every folder setting DEFAULTS to the matching folder of the Eva working
    tree, so a stock C:\Eva install needs no path passed at all. To override one
    (flag beats environment variable beats the CONFIG constant):

        claude mcp add powerpoint --scope user -e PYTHONUTF8=1 -- C:\path\to\python.exe C:\path\to\powerpoint.py --docs-dir D:\Eva\documents\powerpoint --templates-dir D:\Eva\templates\powerpoint --kb-dir D:\Eva\knowledge\powerpoint

    The --docs-dir folder is REQUIRED (flag, POWERPOINT_DOCS_DIR, or the
    DOCS_DIR constant - default C:\Eva\documents\powerpoint): all open/save
    paths must be inside it, NEW decks from powerpoint_create land in it, and
    the server refuses to start without one. It is the ONE folder of
    presentations - there is no separate output folder, so a deck Eva builds
    sits alongside the ones you gave it.
    The --templates-dir folder (POWERPOINT_TEMPLATES_DIR / TEMPLATES_DIR -
    default C:\Eva\templates\powerpoint) holds blank .pptx/.potx templates. It
    is readable like the presentation root but NEVER writable. It is this
    server's own folder: the `word` plugin has its own alongside it, at
    C:\Eva\templates\word.
    The --kb-dir folder (POWERPOINT_KB_DIR / KB_DIR - default
    C:\Eva\knowledge\powerpoint) turns on Markdown mirroring for a local RAG
    knowledge base, on open, create and save; pass 'off' to disable.

    Every one of those folders is part of the Eva working tree: copy the repo's
    eva\ folder to C:\Eva and they all exist, correctly related to each other.
    See eva\README.md for the layout and the reasoning behind it.

    See README.md next to this file for the full settings reference.

=============================================================================
 PROTOCOL / TRANSPORT NOTES
=============================================================================
    - MCP stdio transport = newline-delimited JSON-RPC 2.0 on stdin/stdout.
    - stdout is SACRED: only JSON-RPC messages go there, one per line. Every
      diagnostic goes to stderr via log(). Any stray print() to stdout would
      corrupt the stream.
    - Streams are reconfigured to UTF-8 in-script as a belt-and-braces measure
      alongside PYTHONUTF8.

Author's assumptions (flagged per the airgap "a caveat is cheaper than a
failed transfer" rule):
    - A reasonably modern python-pptx (>= 0.6.21; developed and tested here on
      1.0.2).
    - A .potx template is opened by python-pptx exactly like a .pptx (the part
      names differ, the package does not), so both are accepted as templates.
      A deck is always SAVED as .pptx - saving a .potx back out as a template is
      not supported and the extension is corrected if you ask for one.
    - The speaking-time estimate is word count / SPEAKING_WORDS_PER_MINUTE over
      the SPEAKER NOTES, plus a fixed per-slide allowance. It is a planning aid,
      not a stopwatch: nobody's delivery matches a constant, and a slide with no
      notes contributes only the per-slide allowance. It is reported with its
      inputs so it can be judged.
    - Table cell text is styled by the table style, which is a theme-level
      thing; the font audit reads table text but the fix for an undersized
      table is fewer columns, not a font override.
"""

# Semantic version of this server. Bump on EVERY change (see CLAUDE.md):
# MAJOR = breaking config/tool change, MINOR = new feature, PATCH = fix.
__version__ = "3.0.0"

# =============================================================================
# CONFIGURATION  (all user-editable settings live here, nothing scattered below)
# =============================================================================
SERVER_NAME = "powerpoint"         # advertised to the MCP client
SERVER_VERSION = __version__
PROTOCOL_VERSION_FALLBACK = "2024-11-05"  # used if the client sends none

# REQUIRED path sandbox. The server refuses to open or save any file outside
# this directory tree, and REFUSES TO START if no root is configured - the
# model chooses open/save paths, so an unconfined server could read/write any
# .pptx this account can. Set it here, or at launch with --docs-dir or the
# POWERPOINT_DOCS_DIR environment variable (which take priority over this
# constant). Symlinks are resolved before the containment check.
# This root is ALSO the base for relative paths: a bare "Kickoff.pptx" is
# resolved against it (not the process CWD), so the model can open a file by
# name without knowing its absolute path.
# It is ALSO where powerpoint_create writes new decks: one folder holds every
# .pptx, whether you put it there or Eva built it.
# Default: the Eva working tree's PowerPoint library, searched recursively.
# (--check is exempt: the self-test sandboxes itself to its own temp folder.)
# Related caution: only open .pptx files from trusted sources - a maliciously
# crafted file could use XML entity tricks to pull local file contents into
# the slide text that the model then reads.
DOCS_DIR = r"C:\Eva\documents\powerpoint"

# OPTIONAL folder of blank .pptx/.potx TEMPLATES (the branded deck shell with
# your title slide, section divider and content layouts). Set it here, or at
# launch with --templates-dir or the POWERPOINT_TEMPLATES_DIR environment
# variable (which take priority over this constant). It is a READ-ONLY second
# root:
#   - its files can be listed (powerpoint_list_presentations), opened
#     (powerpoint_open) and used as the base for a new deck
#     (powerpoint_create template="..."), exactly like the presentation root;
#   - every SAVE whose target lands inside it is REFUSED, so a template cannot
#     be overwritten with a filled-in copy of itself. New decks always go to
#     DOCS_DIR.
# One folder per plugin, mirroring DOCS_DIR: the `word` plugin's blanks sit
# beside these in C:\Eva\templates\word. The server also refuses to start if
# this folder IS - or contains - DOCS_DIR, because that arrangement would refuse
# every save.
# Default: the powerpoint\ folder of the Eva working tree's templates zone.
TEMPLATES_DIR = r"C:\Eva\templates\powerpoint"

# OPTIONAL knowledge-base (RAG) folder. If set, EVERY deck opened, created or
# saved is ALSO written out as a Markdown file into this folder, the same way
# word.py mirrors documents, so the content can feed a local RAG index. The
# files are named 'PowerPoint - <name>.md' and overwritten each time. This
# folder is written to by the server only (the model never chooses the path),
# so it does not need to sit inside DOCS_DIR. Set it here, or at launch with
# --kb-dir or the POWERPOINT_KB_DIR environment variable (which take priority
# over this constant). Pass 'off' to disable mirroring.
# Default: the powerpoint\ sub-folder of the Eva knowledge base. It MUST stay
# inside the knowledge-base plugin's documents folder (C:\Eva\knowledge) or the
# mirrored Markdown would never be indexed.
KB_DIR = r"C:\Eva\knowledge\powerpoint"

# --- Guy Kawasaki's 10/20/30 rule ------------------------------------------
# The thresholds powerpoint_review audits against. They are constants rather
# than tool arguments so a house style is set once, in one place, and every
# review in every session applies it. The matching skill
# (skills/kawasaki/SKILL.md) explains the rule; these numbers enforce it.
#   KAWASAKI_MAX_SLIDES  - "10 slides". The pitch-deck ideal.
#   KAWASAKI_MAX_MINUTES - "20 minutes". The speaking budget.
#   KAWASAKI_MIN_FONT_PT - "30-point font". The floor for readable body text.
KAWASAKI_MAX_SLIDES = 10
KAWASAKI_MAX_MINUTES = 20
KAWASAKI_MIN_FONT_PT = 30

# Words per minute used to turn a speaker-note word count into an estimated
# speaking time. 130 is a common figure for measured conference delivery;
# a fast presenter hits 160, a deliberate one 110. Lower it to be pessimistic.
SPEAKING_WORDS_PER_MINUTE = 130

# Seconds added per slide on top of the notes estimate, covering the pause to
# advance, the audience reading the headline, and the sentence you always add
# on the day. Set to 0 to score the notes alone.
SPEAKING_SECONDS_PER_SLIDE = 15

# PowerPoint's own fallback size (points) when nothing in the whole inheritance
# chain specifies one. Used only as the last step of _effective_font_pt.
POWERPOINT_DEFAULT_FONT_PT = 18

# Placeholder types EXEMPT from the 30-point rule. Slide numbers, footers and
# dates are chrome: they are meant to be small, and flagging them would bury
# the real findings. Speaker notes are never audited (nobody reads them from
# the back of the room).
FONT_AUDIT_EXEMPT_PLACEHOLDERS = frozenset((
    "SLIDE_NUMBER", "FOOTER", "DATE",
))

MAX_SESSIONS = 32                    # guard against runaway open() calls
MAX_LAYOUTS_RETURNED = 60            # cap powerpoint_list_layouts output
MAX_FONT_FINDINGS = 80               # cap powerpoint_review's findings list

# powerpoint_add_slide appends ONE slide to the END of the deck, so a deck's
# slide order is simply the order the calls REACH this server. An MCP client is
# free to dispatch independent tool calls concurrently, and concurrent calls do
# not necessarily arrive in the order the model wrote them - which is how a deck
# comes out with its slides shuffled (and why a slide_index guessed rather than
# read back can address the wrong slide). powerpoint_add_slides exists to make
# that impossible: one call, one ordered list of slides. This threshold drives
# the safety net for callers still using the one-slide-per-call tool: when
# several of them land on the same session inside this many seconds - the
# signature of a parallel batch rather than of separate turns - the result
# carries an order_warning telling the caller to verify the order and switch to
# powerpoint_add_slides.
APPEND_BURST_SECONDS = 1.0
APPEND_BURST_MIN_CALLS = 3           # warn from the Nth rapid append onwards

# powerpoint_open name matching. An exact filename always wins; only when no
# exact match is found does the server fall back to a FUZZY match on the name,
# so a near-miss like "quarterly deck" still finds "Quarterly Deck 2024.pptx".
# These tune that fallback (same values as word.py / pdf-to-md.py):
#   FUZZY_MIN_RATIO       - below this similarity, AND with no shared words, a
#                           name is treated as "no match" rather than opened.
#   FUZZY_AMBIGUITY_DELTA - if a runner-up scores within this of the best, the
#                           match is too close to call and the candidates are
#                           listed instead of one being opened silently.
FUZZY_MIN_RATIO = 0.40
FUZZY_AMBIGUITY_DELTA = 0.05

# File extensions this server will open. .potx is PowerPoint's template format;
# python-pptx reads it exactly like a .pptx, so it is accepted as a template.
# Decks are always SAVED as .pptx.
PRESENTATION_EXTENSIONS = (".pptx", ".potx")
# =============================================================================

import sys
import os
import re
import json
import time
import uuid
import difflib
import argparse
import traceback
import datetime

# --- Make stdio UTF-8 regardless of the Windows console codepage. -----------
# Belt-and-braces alongside PYTHONUTF8=1 in the MCP client's env block.
for _stream in ("stdin", "stdout", "stderr"):
    try:
        getattr(sys, _stream).reconfigure(encoding="utf-8")
    except Exception:
        # Older/edge interpreters may not support reconfigure; PYTHONUTF8 covers it.
        pass


def log(msg):
    """All diagnostics go to stderr ONLY. stdout is reserved for JSON-RPC."""
    try:
        sys.stderr.write("[{}] {}\n".format(SERVER_NAME, msg))
        sys.stderr.flush()
    except Exception:
        pass


# --version must work even when python-pptx is not installed yet (useful on an
# endpoint before the wheels are sideloaded), so answer it before the engine
# import below can fail.
if "--version" in sys.argv:
    print("{0} {1}".format(SERVER_NAME, __version__))
    sys.exit(0)

# --- Import the engine, failing loudly and specifically on mismatch. --------
try:
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.exc import PackageNotFoundError
    from pptx.oxml.ns import qn
except Exception as _imp_err:  # pragma: no cover - exercised only on a broken install
    sys.stderr.write(
        "[{}] FATAL: could not import python-pptx.\n"
        "        Interpreter : {}\n"
        "        Error       : {}\n"
        "        Fix: install python-pptx into THIS interpreter "
        "(see the docstring's airgapped install steps).\n".format(
            SERVER_NAME, sys.executable, _imp_err
        )
    )
    sys.exit(1)


def _versions():
    """Best-effort version strings for startup diagnostics."""
    out = {}
    try:
        import pptx as _pptx
        out["python-pptx"] = getattr(_pptx, "__version__", "unknown")
    except Exception:
        out["python-pptx"] = "unknown"
    try:
        from lxml import etree
        out["lxml"] = etree.__version__
    except Exception:
        out["lxml"] = "unknown"
    try:
        import PIL
        out["Pillow"] = getattr(PIL, "__version__", "unknown")
    except Exception:
        out["Pillow"] = "unknown"
    return out


# =============================================================================
# SESSION STATE
# =============================================================================
# session_id -> {"path": str, "prs": Presentation, "opened_at": iso str}
SESSIONS = {}


class ToolError(Exception):
    """Raised by a tool handler to return a clean isError result to the client."""


def _now_iso():
    return datetime.datetime.now().isoformat(timespec="seconds")


def _require(args, key):
    """Fetch a required argument or raise a clear ToolError."""
    if key not in args or args[key] is None:
        raise ToolError("Missing required argument: '{}'".format(key))
    return args[key]


# =============================================================================
# PATH SANDBOX
# =============================================================================
# The model chooses open/save paths, so every one of them is resolved to a real
# path and checked for containment in a permitted root before anything is
# touched. This mirrors word.py exactly - the one write-capable pattern in the
# suite is worth keeping identical across servers so a reader only learns it
# once.
def _permitted_roots():
    """
    Folders the server is allowed to READ inside: the DOCS_DIR sandbox, plus the
    TEMPLATES_DIR of blank templates when one is configured. The knowledge-base
    folder is deliberately NOT included - the model never supplies a path into
    it; the server writes there itself.

    Order matters: a relative name that could live in either root is resolved
    against these in turn, and a not-yet-existing file (a save-as target) falls
    back to the FIRST one, which must stay the presentation root. Writes are
    further restricted - see _read_only_root().
    """
    roots = []
    if DOCS_DIR:
        roots.append(DOCS_DIR)
    if TEMPLATES_DIR:
        roots.append(TEMPLATES_DIR)
    return roots


def _read_only_root(rp):
    """
    Return the READ-ONLY root containing resolved path `rp`, or None.

    Only TEMPLATES_DIR is read-only: templates are a library of blanks to start
    from, so they are readable everywhere a deck is, but nothing may be saved
    over them.
    """
    if not TEMPLATES_DIR:
        return None
    real_root = os.path.realpath(os.path.expanduser(TEMPLATES_DIR))
    try:
        # Different drives on Windows raise ValueError from commonpath.
        if os.path.commonpath([real_root, rp]) == real_root:
            return real_root
    except ValueError:
        return None
    return None


def _refuse_if_read_only(rp, what="Saving"):
    """Guard every write: refuse a target inside the read-only templates folder
    and say where the file should go instead."""
    root = _read_only_root(rp)
    if root is not None:
        raise ToolError(
            "{} into the templates folder is not allowed - it is read-only so "
            "templates stay blank ({}). Use powerpoint_create (which writes to "
            "the presentations folder) to start a new deck from a template, or "
            "save-as to a path inside the presentations folder.".format(what, root)
        )


def _root_label(rp):
    """
    Which configured folder a resolved path lives in: 'docs' or 'templates'
    (None if somehow outside them both). The MOST SPECIFIC root wins, so a
    templates folder nested inside the presentation root is still reported as
    'templates'.
    """
    label, best_len = None, -1
    for name, root in (("docs", DOCS_DIR), ("templates", TEMPLATES_DIR)):
        if not root:
            continue
        real_root = os.path.realpath(os.path.expanduser(root))
        try:
            if os.path.commonpath([real_root, rp]) == real_root and \
                    len(real_root) > best_len:
                label, best_len = name, len(real_root)
        except ValueError:
            continue
    return label


def _contained_in_root(rp):
    """
    Return the permitted root that contains resolved path `rp`, or None. A
    symlink is already resolved by the caller (realpath), so a link inside a
    root cannot smuggle the target outside it.
    """
    for root in _permitted_roots():
        real_root = os.path.realpath(os.path.expanduser(root))
        try:
            # Different drives on Windows raise ValueError from commonpath.
            if os.path.commonpath([real_root, rp]) == real_root:
                return real_root
        except ValueError:
            continue
    return None


def _resolve_path(path):
    """
    Turn a caller-supplied path into an absolute, sandbox-checked path.

    A RELATIVE path (e.g. just "Kickoff.pptx") is resolved against the permitted
    roots - the presentation root first, then the templates folder - NOT against
    the server process's current working directory. The CWD is wherever the MCP
    client launched python and is almost never the presentation folder. When a
    relative name maps into both roots, an existing file is preferred; otherwise
    the docs-dir candidate wins (the natural target for a new save-as).

    An ABSOLUTE (or ~) path is taken as-is and must still fall inside a
    permitted root. realpath is used throughout so symlinks are resolved before
    the containment check.
    """
    if not isinstance(path, str) or not path.strip():
        raise ToolError("Path must be a non-empty string")
    roots = _permitted_roots()
    if not roots:
        # main() refuses to start without a root; this guards direct callers.
        raise ToolError("No DOCS_DIR is configured; file access is disabled.")

    expanded = os.path.expanduser(path.strip())
    if os.path.isabs(expanded):
        candidates = [os.path.realpath(expanded)]
    else:
        # Interpret the relative name inside each permitted root.
        candidates = [
            os.path.realpath(os.path.join(os.path.realpath(os.path.expanduser(root)),
                                          expanded))
            for root in roots
        ]

    contained = [rp for rp in candidates if _contained_in_root(rp) is not None]
    if not contained:
        raise ToolError(
            "Path is outside the permitted folder(s) (DOCS_DIR"
            + (" / TEMPLATES_DIR" if TEMPLATES_DIR else "")
            + ") and was refused."
        )
    # Prefer a candidate that already exists (matters when a relative name could
    # live in either root); else fall back to the first (the docs dir).
    for rp in contained:
        if os.path.isfile(rp):
            return rp
    return contained[0]


def _iter_pptx_files(root):
    """Yield every .pptx/.potx file under `root` (recursively), skipping Office
    lock files like '~$Kickoff.pptx'."""
    for dirpath, _dirnames, filenames in os.walk(root):
        for fn in filenames:
            if fn.lower().endswith(PRESENTATION_EXTENSIONS) and not fn.startswith("~$"):
                yield os.path.join(dirpath, fn)


def _all_pptx_files():
    """All presentations under the permitted roots, de-duplicated, in root
    order."""
    seen = set()
    out = []
    for root in _permitted_roots():
        real_root = os.path.realpath(os.path.expanduser(root))
        if not os.path.isdir(real_root):
            continue
        for p in _iter_pptx_files(real_root):
            rp = os.path.realpath(p)
            if rp not in seen:
                seen.add(rp)
                out.append(rp)
    return out


def _display_path(p):
    """Path shown to the model: relative to its containing root when possible
    (so it can be handed straight back to powerpoint_open), else absolute."""
    root = _contained_in_root(p)
    if root is not None:
        try:
            return os.path.relpath(p, root)
        except ValueError:
            pass
    return p


def _find_pptx_by_name(name):
    """
    Full paths under the permitted roots whose FILENAME matches `name`
    (case-insensitive). A name with no extension matches either supported
    extension. Any directory part of `name` is ignored - this is a last-resort
    basename lookup so a bare filename opens even when it sits in a subfolder.
    """
    want = os.path.basename(name.strip().replace("\\", "/")).lower()
    if not want:
        return []
    if want.endswith(PRESENTATION_EXTENSIONS):
        wanted = (want,)
    else:
        wanted = tuple(want + ext for ext in PRESENTATION_EXTENSIONS)
    return [p for p in _all_pptx_files() if os.path.basename(p).lower() in wanted]


def _normalise_name(text):
    """Lowercase, turn separators into spaces and collapse whitespace, so fuzzy
    matching ignores punctuation/case differences between a query and a name."""
    text = text.lower()
    for ch in ("_", "-", ".", "(", ")", "[", "]", ",", "&", "+"):
        text = text.replace(ch, " ")
    return " ".join(text.split())


def _score_name(query_norm, name_norm):
    """Score a query against a name as (token-containment, sequence-ratio): the
    fraction of query words present in the name, then difflib's overall
    similarity. Ranking on the pair prefers names that contain the query words
    and, among those, the closest overall string."""
    q_tokens = query_norm.split()
    n_tokens = set(name_norm.split())
    contained = (sum(1 for t in q_tokens if t in n_tokens) / len(q_tokens)) if q_tokens else 0.0
    ratio = difflib.SequenceMatcher(None, query_norm, name_norm).ratio()
    return (contained, ratio)


def _fuzzy_match_pptx(name):
    """
    Fuzzy-match `name` against the presentations under the permitted roots.
    Returns (best_path, tied): best_path is the single best match, or None when
    nothing is close enough; `tied` lists the near-equal candidates when the
    match is too ambiguous to pick one. Matching is on the file's stem.
    """
    files = _all_pptx_files()
    if not files:
        return None, []
    query_stem = os.path.splitext(os.path.basename(name.strip().replace("\\", "/")))[0]
    query_norm = _normalise_name(query_stem)
    if not query_norm:
        return None, []
    scored = [
        (p, _score_name(query_norm, _normalise_name(os.path.splitext(os.path.basename(p))[0])))
        for p in files
    ]
    scored.sort(key=lambda item: item[1], reverse=True)
    best_path, (best_contained, best_ratio) = scored[0]
    # Nothing shares a word and the closest name is still too different.
    if best_contained == 0.0 and best_ratio < FUZZY_MIN_RATIO:
        return None, []
    tied = []
    for p, (contained, ratio) in scored:
        if contained == best_contained and (best_ratio - ratio) <= FUZZY_AMBIGUITY_DELTA:
            tied.append(p)
        else:
            break
    if len(tied) > 1:
        return None, tied
    return best_path, []


def _resolve_existing_pptx(raw):
    """
    Resolve a caller-supplied name/path to an EXISTING presentation under the
    permitted roots, using the forgiving chain shared by powerpoint_open and the
    template option of powerpoint_create: sandbox-resolve the path, then (if not
    found) an exact basename lookup, then a fuzzy name match. Returns
    (path, fuzzy_matched). Raises ToolError if nothing matches, the name is
    ambiguous, or the resolved file is not a presentation.
    """
    path = _resolve_path(raw)
    fuzzy_matched = False
    if not os.path.isfile(path):
        # Not at the resolved location. Fall back to a basename lookup anywhere
        # under the permitted roots, so "Deck Template.pptx" resolves even when
        # it lives in a subfolder of the presentation root.
        matches = _find_pptx_by_name(raw)
        if len(matches) == 1:
            path = matches[0]
        elif len(matches) > 1:
            listing = ", ".join(_display_path(m) for m in matches)
            raise ToolError(
                "Several presentations are named '{}': {}. Use the path "
                "(relative to the presentation root) to pick one.".format(
                    os.path.basename(raw.replace("\\", "/")), listing)
            )
        else:
            # No exact filename match - fall back to a fuzzy match on the name.
            best, tied = _fuzzy_match_pptx(raw)
            if best is not None:
                log("fuzzy-matched '{}' -> {}".format(raw, best))
                path = best
                fuzzy_matched = True
            elif tied:
                listing = ", ".join(_display_path(m) for m in tied)
                raise ToolError(
                    "'{}' matches several presentations about equally well: {}. "
                    "Use the exact name or path, or call "
                    "powerpoint_list_presentations.".format(raw, listing)
                )
            else:
                raise ToolError(
                    "File not found: '{}'. Paths are resolved relative to the "
                    "presentation root; call powerpoint_list_presentations to "
                    "see what is available.".format(raw)
                )
    if not path.lower().endswith(PRESENTATION_EXTENSIONS):
        raise ToolError(
            "Only .pptx and .potx files are supported (got: {})".format(path)
        )
    return path, fuzzy_matched


def _get_session(args):
    """Return the session dict for a required session_id argument."""
    sid = _require(args, "session_id")
    session = SESSIONS.get(sid)
    if session is None:
        raise ToolError(
            "Unknown session_id '{}'. Call powerpoint_create or "
            "powerpoint_open first.".format(sid)
        )
    return session


def safe_filename(name, max_len=150):
    """Strip characters Windows forbids in a filename, collapse whitespace and
    trim to a sane length, so a slide title can be used as a file name."""
    cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "", str(name)).strip()
    cleaned = re.sub(r"\s+", " ", cleaned).strip(". ")
    if not cleaned:
        cleaned = "untitled"
    return cleaned[:max_len]


# =============================================================================
# TEMPLATE INTROSPECTION  (layouts, placeholders, and the font chain)
# =============================================================================
# Everything in this section exists to answer one question the naive python-pptx
# API cannot: "what will this text ACTUALLY look like when the template is
# applied?". A deck that inherits properly has almost no explicit formatting, so
# `run.font.size` is None nearly everywhere and the interesting values live in
# the layout, the master and the theme.

# Namespaced tag names, resolved once. qn() maps the prefixes python-pptx
# registers: 'a:' = DrawingML, 'p:' = PresentationML.
_A_LSTSTYLE = qn("a:lstStyle")
_A_PPR = qn("a:pPr")
_A_RPR = qn("a:rPr")
_A_DEFRPR = qn("a:defRPr")
_P_TXSTYLES = qn("p:txStyles")
_P_TITLESTYLE = qn("p:titleStyle")
_P_BODYSTYLE = qn("p:bodyStyle")
_P_OTHERSTYLE = qn("p:otherStyle")
_P_DEFAULTTEXTSTYLE = qn("p:defaultTextStyle")
_P_SLDIDLST = qn("p:sldIdLst")
_P_SLDID = qn("p:sldId")

# Placeholder types that carry the deck's real content, as opposed to the
# date/footer/slide-number chrome. Used when deciding which unfilled
# placeholders to drop and which text to audit.
_CONTENT_PH_TYPES = frozenset((
    "TITLE", "CENTER_TITLE", "SUBTITLE", "BODY", "OBJECT",
    "TABLE", "CHART", "PICTURE", "CLIP_ART", "MEDIA_CLIP",
    "ORG_CHART", "BITMAP",
))

# Which of the master's three txStyles blocks governs a given placeholder type.
# Anything not listed (and every non-placeholder shape) falls to otherStyle,
# which is what PowerPoint itself does.
_TITLE_PH_TYPES = frozenset(("TITLE", "CENTER_TITLE"))
_BODY_PH_TYPES = frozenset((
    "BODY", "SUBTITLE", "OBJECT", "TABLE", "CHART", "ORG_CHART", "MEDIA_CLIP",
))


def _ph_type_name(shape):
    """
    The placeholder type of `shape` as a plain string ('TITLE', 'BODY', ...),
    or None if the shape is not a placeholder.

    python-pptx returns an enum member whose str() is like 'BODY (2)'; the bare
    name is what the tool results and the config constants speak in.
    """
    try:
        if not shape.is_placeholder:
            return None
        ph_type = shape.placeholder_format.type
    except (AttributeError, ValueError, KeyError):
        return None
    if ph_type is None:
        return None
    # Enum members expose .name; fall back to parsing str() for older releases.
    name = getattr(ph_type, "name", None)
    if name:
        return str(name)
    return str(ph_type).split(" ")[0]


def _ph_idx(shape):
    """The placeholder index of `shape`, or None if it is not a placeholder.
    idx is what ties a slide placeholder to the layout placeholder it inherits
    from - the type alone is ambiguous (a Two Content layout has two OBJECTs)."""
    try:
        if not shape.is_placeholder:
            return None
        return shape.placeholder_format.idx
    except (AttributeError, ValueError, KeyError):
        return None


def _find_placeholder_by_idx(container, idx):
    """The placeholder on `container` (a slide, layout or master) whose idx
    matches, or None. Used to walk one step up the inheritance chain."""
    if idx is None:
        return None
    try:
        for ph in container.placeholders:
            if _ph_idx(ph) == idx:
                return ph
    except (AttributeError, KeyError):
        return None
    return None


def _find_master_placeholder_by_type(master, ph_type_name):
    """
    The master placeholder that governs `ph_type_name`, or None.

    Masters do not use the layouts' idx numbering: they carry one placeholder
    per ROLE (title, body, date, footer, slide number), so a layout's
    CENTER_TITLE inherits from the master's TITLE and every body-ish type
    inherits from the master's BODY.
    """
    if master is None or not ph_type_name:
        return None
    if ph_type_name in _TITLE_PH_TYPES:
        wanted = _TITLE_PH_TYPES
    elif ph_type_name in _BODY_PH_TYPES:
        wanted = _BODY_PH_TYPES
    else:
        wanted = frozenset((ph_type_name,))
    try:
        for ph in master.placeholders:
            if _ph_type_name(ph) in wanted:
                return ph
    except (AttributeError, KeyError):
        return None
    return None


def _sz_from_lvl_container(container_el, level):
    """
    Read the font size (in points, as a float) for outline `level` (0-based)
    out of an element that holds a:lvl1pPr..a:lvl9pPr children - an a:lstStyle,
    a p:titleStyle/p:bodyStyle/p:otherStyle, or a p:defaultTextStyle.

    Returns None when that element does not specify a size for the level.
    PowerPoint stores sizes in hundredths of a point, so sz="3200" is 32pt.
    """
    if container_el is None:
        return None
    # Levels beyond 9 reuse level 9, which is what PowerPoint does.
    lvl = min(max(int(level), 0), 8) + 1
    lvl_el = container_el.find(qn("a:lvl{}pPr".format(lvl)))
    if lvl_el is None:
        return None
    def_rpr = lvl_el.find(_A_DEFRPR)
    if def_rpr is None:
        return None
    sz = def_rpr.get("sz")
    if not sz:
        return None
    try:
        return int(sz) / 100.0
    except (TypeError, ValueError):
        return None


def _sz_from_shape_lst_style(shape, level):
    """The size a shape's own a:lstStyle sets for `level`, or None. This is the
    per-shape override a designer applies when one box on one slide differs from
    the layout."""
    try:
        txbody = shape.text_frame._txBody
    except (AttributeError, ValueError):
        return None
    return _sz_from_lvl_container(txbody.find(_A_LSTSTYLE), level)


def _sz_from_master_txstyles(master, ph_type_name, level):
    """
    The size the master's p:txStyles block sets for `level`, or None.

    This is the step that actually supplies the number in most real templates:
    the stock Office master sets titleStyle lvl1 to 44pt and bodyStyle lvl1/2/3
    to 32/28/24pt, and almost every branded template does the same thing with
    its own numbers. It is also why a deck can fail the 30-point rule without a
    single explicit font size anywhere in it.
    """
    if master is None:
        return None
    tx_styles = master.element.find(_P_TXSTYLES)
    if tx_styles is None:
        return None
    if ph_type_name in _TITLE_PH_TYPES:
        block = tx_styles.find(_P_TITLESTYLE)
    elif ph_type_name in _BODY_PH_TYPES:
        block = tx_styles.find(_P_BODYSTYLE)
    else:
        # Non-placeholder shapes and chrome placeholders use otherStyle.
        block = tx_styles.find(_P_OTHERSTYLE)
    return _sz_from_lvl_container(block, level)


def _sz_from_presentation_default(prs, level):
    """The size the presentation's p:defaultTextStyle sets for `level`, or None.
    The last stop before PowerPoint's hard-coded fallback."""
    try:
        return _sz_from_lvl_container(
            prs.part._element.find(_P_DEFAULTTEXTSTYLE), level)
    except (AttributeError, ValueError):
        return None


def _chain_for(container):
    """
    The (layout, master) a shape's formatting inherits from, given the thing the
    shape sits on. Works for all three: a SLIDE inherits from its layout and
    that layout's master; a LAYOUT inherits from its master only (its own
    placeholders are already step 3 of the chain); a MASTER inherits from
    nothing above itself.

    Getting this right is what lets powerpoint_list_layouts report a template's
    real font sizes BEFORE any slide exists - the common case of "will this
    template hold 30-point text?".
    """
    layout = getattr(container, "slide_layout", None)
    if layout is not None:
        return layout, getattr(layout, "slide_master", None)   # container is a slide
    master = getattr(container, "slide_master", None)
    if master is not None:
        return None, master                                    # container is a layout
    return None, container                                     # container is a master


def _effective_font_pt(prs, container, shape, paragraph, run, level):
    """
    Resolve the size text will ACTUALLY render at, in points, the way PowerPoint
    resolves it. `container` is the slide, layout or master the shape sits on.
    Returns (points, source) where `source` names the step that supplied the
    value, so a finding can explain itself:

        run              the run's own a:rPr/@sz
        paragraph        the paragraph's a:pPr/a:defRPr/@sz
        shape            the shape's a:lstStyle for this outline level
        layout           the matching placeholder on the slide LAYOUT
        master           the matching placeholder on the slide MASTER
        master_txstyles  the master's titleStyle / bodyStyle / otherStyle
        presentation     the presentation's p:defaultTextStyle
        fallback         POWERPOINT_DEFAULT_FONT_PT - nothing specified one

    `run` may be None to ask what an EMPTY paragraph (or a whole placeholder)
    would render at, which is how powerpoint_list_layouts reports a layout's
    sizes before any text exists.
    """
    # 1. The run's own explicit size.
    if run is not None:
        try:
            if run.font.size is not None:
                return (run.font.size.pt, "run")
        except (AttributeError, ValueError):
            pass

    # 2. The paragraph's default run properties.
    if paragraph is not None:
        try:
            p_el = paragraph._p
            ppr = p_el.find(_A_PPR)
            if ppr is not None:
                def_rpr = ppr.find(_A_DEFRPR)
                if def_rpr is not None and def_rpr.get("sz"):
                    return (int(def_rpr.get("sz")) / 100.0, "paragraph")
        except (AttributeError, ValueError, TypeError):
            pass

    # 3. The shape's own list style.
    pt = _sz_from_shape_lst_style(shape, level)
    if pt is not None:
        return (pt, "shape")

    ph_type_name = _ph_type_name(shape)
    ph_idx = _ph_idx(shape)
    layout, master = _chain_for(container)

    # 4. The matching placeholder on the layout (by idx - the type is ambiguous
    #    when a layout has two content placeholders). Skipped when the shape IS
    #    a layout placeholder, since step 3 already read it.
    layout_ph = _find_placeholder_by_idx(layout, ph_idx) if layout is not None else None
    if layout_ph is not None:
        pt = _sz_from_shape_lst_style(layout_ph, level)
        if pt is not None:
            return (pt, "layout")

    # 5. The matching placeholder on the master (by role).
    master_ph = _find_master_placeholder_by_type(master, ph_type_name)
    if master_ph is not None:
        pt = _sz_from_shape_lst_style(master_ph, level)
        if pt is not None:
            return (pt, "master")

    # 6. The master's txStyles - where most templates really set their sizes.
    pt = _sz_from_master_txstyles(master, ph_type_name, level)
    if pt is not None:
        return (pt, "master_txstyles")

    # 7. The presentation-wide default text style.
    pt = _sz_from_presentation_default(prs, level)
    if pt is not None:
        return (pt, "presentation")

    # 8. PowerPoint's own fallback.
    return (float(POWERPOINT_DEFAULT_FONT_PT), "fallback")


def _layout_placeholder_summary(prs, layout, ph):
    """One entry describing a placeholder on a LAYOUT: how to address it, and
    what size level-1 text dropped into it would come out at. The size is the
    whole point - it is how you tell, before writing a word, whether a layout
    can hold 30-point text."""
    ph_type_name = _ph_type_name(ph) or "UNKNOWN"
    # A layout placeholder resolves against its own layout and master; passing
    # the layout as the "slide" walks the same chain from one step higher up.
    pt, source = _effective_font_pt(prs, layout, ph, None, None, 0)
    entry = {
        "idx": _ph_idx(ph),
        "type": ph_type_name,
        "name": ph.name,
        "level1_font_pt": round(pt, 1),
        "level1_font_source": source,
    }
    if ph_type_name in _BODY_PH_TYPES:
        # Bulleted content is the usual 30-point offender, and it is the deeper
        # levels that break it, so report them too.
        levels = []
        for lvl in range(1, 4):
            lvl_pt, lvl_src = _effective_font_pt(prs, layout, ph, None, None, lvl)
            levels.append({"level": lvl + 1,
                           "font_pt": round(lvl_pt, 1),
                           "source": lvl_src})
        entry["deeper_levels"] = levels
    return entry


def _layout_role(layout):
    """
    Classify a layout by what its placeholders are, NOT by its name - a branded
    template calls its layouts whatever it likes ('Divider', 'Chapter Opener',
    'Standard Content'), so matching on names is a guessing game. Returns one
    of: 'title', 'section', 'bullets', 'two_content', 'title_only', 'picture',
    'blank', 'other'.
    """
    types = [_ph_type_name(ph) for ph in layout.placeholders]
    content = [t for t in types if t in _CONTENT_PH_TYPES]
    has_ctr_title = "CENTER_TITLE" in types
    has_title = "TITLE" in types or has_ctr_title
    has_subtitle = "SUBTITLE" in types
    body_like = [t for t in content if t in ("BODY", "OBJECT")]
    has_picture = "PICTURE" in types or "BITMAP" in types

    if not content:
        return "blank"
    if has_ctr_title or (has_title and has_subtitle):
        return "title"
    if has_picture:
        return "picture"
    if has_title and len(body_like) >= 2:
        return "two_content"
    if has_title and len(body_like) == 1:
        # A section divider and a bullets slide look alike structurally; the
        # divider's body is a short standfirst, so fall back to the name for
        # this ONE distinction, where the convention is near-universal.
        if re.search(r"\b(section|divider|chapter|part|header)\b",
                     layout.name or "", re.IGNORECASE):
            return "section"
        return "bullets"
    if has_title and not body_like:
        return "title_only"
    return "other"


def _recommended_layouts(prs):
    """
    The best layout in THIS template for each common job, as
    {role: {"index": int, "name": str}}. Mirrors word.py's list_styles
    `recommended` block: when the user just says "add a slide", this is what
    decides which layout gets used, so the answer comes from the template rather
    than from a hard-coded index that only happens to be right for the stock
    Office deck.
    """
    picked = {}
    for i, layout in enumerate(prs.slide_layouts):
        role = _layout_role(layout)
        if role in ("other",):
            continue
        # First layout of each role wins - templates are ordered deliberately,
        # with the intended one first.
        picked.setdefault(role, {"index": i, "name": layout.name})
    # 'bullets' is the workhorse; if the template has none, a two-content or
    # title-only layout is the least-bad substitute.
    if "bullets" not in picked:
        for fallback in ("two_content", "title_only", "section"):
            if fallback in picked:
                picked["bullets"] = dict(picked[fallback])
                break
    return picked


def _resolve_layout(prs, spec):
    """
    Turn a caller's layout choice into a real layout object. Accepts an integer
    index, a numeric string, an exact/case-insensitive layout NAME, or one of
    the role words from _recommended_layouts ('title', 'section', 'bullets',
    'two_content', 'title_only', 'picture', 'blank'). Raises a ToolError that
    lists the real names when nothing matches - a dead end here is the most
    likely thing to stall a build.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ToolError(
            "This template defines no slide layouts, so no slide can be added.")

    def _listing():
        return ", ".join("{}='{}'".format(i, lay.name) for i, lay in enumerate(layouts))

    if spec is None or (isinstance(spec, str) and not spec.strip()):
        # No preference: use the template's own bullets/content layout.
        rec = _recommended_layouts(prs)
        choice = rec.get("bullets") or rec.get("title_only") or rec.get("title")
        if choice is None:
            return layouts[0]
        return layouts[choice["index"]]

    if isinstance(spec, bool):
        raise ToolError("layout must be a layout index, name or role - not a boolean")

    if isinstance(spec, int):
        if 0 <= spec < len(layouts):
            return layouts[spec]
        raise ToolError(
            "layout index {} is out of range - this template has {} layouts: "
            "{}".format(spec, len(layouts), _listing()))

    text = str(spec).strip()
    if text.isdigit():
        return _resolve_layout(prs, int(text))

    # An exact (then case-insensitive) name match.
    for layout in layouts:
        if layout.name == text:
            return layout
    lowered = text.lower()
    for layout in layouts:
        if (layout.name or "").lower() == lowered:
            return layout

    # A role word.
    role_key = lowered.replace(" ", "_").replace("-", "_")
    rec = _recommended_layouts(prs)
    if role_key in rec:
        return layouts[rec[role_key]["index"]]

    # Last resort: a fuzzy name match, so 'content' finds 'Title and Content'.
    best, tied = None, []
    scored = sorted(
        ((lay, _score_name(_normalise_name(text), _normalise_name(lay.name or "")))
         for lay in layouts),
        key=lambda item: item[1], reverse=True)
    if scored:
        best, (contained, ratio) = scored[0]
        if contained == 0.0 and ratio < FUZZY_MIN_RATIO:
            best = None
        else:
            tied = [lay for lay, (c, r) in scored
                    if c == contained and (ratio - r) <= FUZZY_AMBIGUITY_DELTA]
    if best is not None and len(tied) <= 1:
        log("fuzzy-matched layout '{}' -> '{}'".format(text, best.name))
        return best

    raise ToolError(
        "No layout matches '{}'. This template's layouts are: {}. You can also "
        "pass a role: {}. Call powerpoint_list_layouts to see each layout's "
        "placeholders and font sizes.".format(
            text, _listing(), ", ".join(sorted(rec)) or "(none detected)")
    )


# =============================================================================
# SLIDE CONTENT HELPERS
# =============================================================================
def _shape_text(shape):
    """All text in a shape, paragraphs joined by newlines. '' for a shape with
    no text frame."""
    try:
        if not shape.has_text_frame:
            return ""
        return "\n".join(p.text for p in shape.text_frame.paragraphs)
    except (AttributeError, ValueError):
        return ""


def _is_empty_text_shape(shape):
    """True when a shape has a text frame carrying nothing but whitespace."""
    try:
        if not shape.has_text_frame:
            return False
    except (AttributeError, ValueError):
        return False
    return not _shape_text(shape).strip()


def _slide_title(slide):
    """The slide's title text, or '' - reading the title placeholder if there is
    one, else the first title-typed placeholder that carries text."""
    try:
        if slide.shapes.title is not None:
            text = _shape_text(slide.shapes.title).strip()
            if text:
                return text
    except (AttributeError, ValueError):
        pass
    for shape in slide.shapes:
        if _ph_type_name(shape) in _TITLE_PH_TYPES:
            text = _shape_text(shape).strip()
            if text:
                return text
    return ""


def _notes_text(slide):
    """The slide's speaker notes, or ''. Asking for a notes slide CREATES one in
    python-pptx, so has_notes_slide is checked first - otherwise merely reading
    a deck would dirty every slide in it."""
    try:
        if not slide.has_notes_slide:
            return ""
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except (AttributeError, ValueError):
        return ""


def _iter_text_runs(slide):
    """
    Yield (shape, paragraph, run, level, in_table) for every run of text on a
    slide, including inside table cells. Grouped shapes are walked recursively,
    since a branded template often groups a caption with its rule.
    """
    def _walk(shapes):
        for shape in shapes:
            # Grouped shapes: recurse into the group's members.
            if hasattr(shape, "shapes"):
                for item in _walk(shape.shapes):
                    yield item
                continue
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            yield (shape, para, run, para.level, False)
            except (AttributeError, ValueError):
                pass
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    yield (shape, para, run, para.level, True)
            except (AttributeError, ValueError):
                pass

    for item in _walk(slide.shapes):
        yield item


_TYPED_BULLET_RE = re.compile(r"^\s*(?:[-*•–—]|\d+[.)])\s+")


def _strip_typed_bullet(text):
    """Remove a hand-typed bullet marker ('- ', '* ', '1. ', a bullet glyph)
    from the start of a line. The layout supplies the bullet character, so a
    typed one renders as a doubled marker - the PowerPoint equivalent of the
    fake-list problem word.py guards against."""
    return _TYPED_BULLET_RE.sub("", text, count=1)


def _normalise_bullets(raw, argname="bullets"):
    """
    Accept the several shapes a caller might send bullets in and return a list
    of {"text": str, "level": int}:
        ["a", "b"]                                  -> both at level 0
        [{"text": "a"}, {"text": "b", "level": 1}]  -> explicit levels
        "a\nb"                                      -> split on newlines
    Levels are 0-based (0 = top level), and PowerPoint allows 0-8.
    """
    if raw is None:
        return []
    if isinstance(raw, str):
        raw = [line for line in raw.split("\n") if line.strip()]
    if not isinstance(raw, list):
        raise ToolError(
            "{} must be a list of strings, a list of {{text, level}} objects, "
            "or a newline-separated string.".format(argname))
    items = []
    for i, entry in enumerate(raw):
        if isinstance(entry, str):
            text, level = entry, 0
        elif isinstance(entry, dict):
            if "text" not in entry:
                raise ToolError(
                    "{}[{}] is an object with no 'text' key.".format(argname, i))
            text = entry.get("text")
            level = entry.get("level", 0)
        else:
            raise ToolError(
                "{}[{}] must be a string or an object with 'text' (and "
                "optionally 'level').".format(argname, i))
        if not isinstance(text, str):
            raise ToolError("{}[{}].text must be a string.".format(argname, i))
        try:
            level = int(level or 0)
        except (TypeError, ValueError):
            raise ToolError(
                "{}[{}].level must be a whole number 0-8.".format(argname, i))
        if not 0 <= level <= 8:
            raise ToolError(
                "{}[{}].level is {} - PowerPoint outline levels are 0-8 "
                "(0 = top level).".format(argname, i, level))
        # A leading '- ' or bullet glyph is the template's job, not the text's:
        # a typed one renders on top of the layout's own bullet character.
        items.append({"text": _strip_typed_bullet(text), "level": level})
    return items


def _fill_text_frame(text_frame, items, replace=True):
    """
    Write `items` (from _normalise_bullets) into a text frame, one paragraph per
    item, at the outline level each one asks for. Returns the number written.

    NOTHING here sets a font, size, colour or alignment: the paragraph's level
    plus the placeholder it lives in are what the template styles against, and
    setting anything explicitly is exactly how a deck stops matching its
    template. That restraint is the feature, not an omission.
    """
    if replace:
        # clear() leaves exactly one empty paragraph behind.
        text_frame.clear()
    written = 0
    for item in items:
        para = None
        if written == 0:
            # Reuse a trailing empty paragraph (a freshly cleared frame, or an
            # untouched placeholder) instead of leaving a blank line above the
            # first bullet.
            last = text_frame.paragraphs[-1] if text_frame.paragraphs else None
            if last is not None and not last.text.strip() and not last.runs:
                para = last
        if para is None:
            para = text_frame.add_paragraph()
        para.text = item["text"]
        para.level = item["level"]
        written += 1
    return written


def _drop_empty_content_placeholders(slide):
    """
    Remove content placeholders left empty, so the deck carries no
    "Click to add text" prompts. Returns the names removed.

    Deliberately conservative about what it touches:
      - date / footer / slide-number placeholders are NEVER removed; they are
        chrome the master supplies on purpose;
      - picture / table / chart placeholders are NEVER removed either, because
        an empty one is the click target a human needs to finish the slide by
        hand - this server does not insert images.
    Only empty TEXT placeholders go.
    """
    removed = []
    for shape in list(slide.shapes):
        ph_type = _ph_type_name(shape)
        if ph_type is None or ph_type not in _CONTENT_PH_TYPES:
            continue
        if ph_type in ("PICTURE", "TABLE", "CHART", "CLIP_ART", "BITMAP",
                       "ORG_CHART", "MEDIA_CLIP"):
            continue
        if _is_empty_text_shape(shape):
            shape._element.getparent().remove(shape._element)
            removed.append(shape.name)
    return removed


def _placeholder_lookup(slide, spec):
    """
    Find a placeholder on `slide` by index (its idx), by exact/case-insensitive
    NAME, or by type word ('title', 'body', 'subtitle'). Raises a ToolError
    listing what the slide actually has, so a wrong guess is one call from being
    right.
    """
    available = [
        {"idx": _ph_idx(ph), "type": _ph_type_name(ph), "name": ph.name}
        for ph in slide.placeholders
    ]

    def _fail():
        listing = ", ".join(
            "idx={} type={} name='{}'".format(a["idx"], a["type"], a["name"])
            for a in available) or "(none)"
        raise ToolError(
            "No placeholder matches {!r} on this slide. It has: {}. Note that "
            "powerpoint_add_slide drops EMPTY text placeholders by default, so "
            "fill one when you add the slide, or pass "
            "drop_empty_placeholders=false.".format(spec, listing))

    if spec is None:
        _fail()
    if isinstance(spec, bool):
        raise ToolError("placeholder must be an index, name or type word")
    if isinstance(spec, int) or (isinstance(spec, str) and spec.strip().isdigit()):
        idx = int(spec)
        for ph in slide.placeholders:
            if _ph_idx(ph) == idx:
                return ph
        _fail()

    text = str(spec).strip()
    for ph in slide.placeholders:
        if ph.name == text:
            return ph
    lowered = text.lower()
    for ph in slide.placeholders:
        if (ph.name or "").lower() == lowered:
            return ph
    # A type word: 'title' also accepts the layout's CENTER_TITLE, and 'body'
    # accepts OBJECT, because that distinction is a layout detail the caller
    # should not have to know.
    wanted = lowered.replace(" ", "_").upper()
    aliases = {
        "TITLE": _TITLE_PH_TYPES,
        "HEADING": _TITLE_PH_TYPES,
        "BODY": frozenset(("BODY", "OBJECT")),
        "CONTENT": frozenset(("BODY", "OBJECT")),
        "BULLETS": frozenset(("BODY", "OBJECT")),
        "SUBTITLE": frozenset(("SUBTITLE",)),
    }
    wanted_set = aliases.get(wanted, frozenset((wanted,)))
    for ph in slide.placeholders:
        if _ph_type_name(ph) in wanted_set:
            return ph
    _fail()


def _get_slide(prs, index, argname="slide_index"):
    """Fetch a slide by zero-based index, with a message that says what the
    valid range actually is."""
    if isinstance(index, bool) or not isinstance(index, int):
        try:
            index = int(index)
        except (TypeError, ValueError):
            raise ToolError("{} must be a whole number.".format(argname))
    slides = list(prs.slides)
    if not slides:
        raise ToolError("This deck has no slides yet - add one with "
                        "powerpoint_add_slide.")
    if not 0 <= index < len(slides):
        raise ToolError(
            "{} {} is out of range - this deck has {} slides (0-{}).".format(
                argname, index, len(slides), len(slides) - 1))
    return slides[index]


def _table_rows(shape):
    """A shape's table as a list of row lists of cell text, or None."""
    try:
        if not shape.has_table:
            return None
    except (AttributeError, ValueError):
        return None
    return [[cell.text for cell in row.cells] for row in shape.table.rows]


# =============================================================================
# KNOWLEDGE-BASE MIRROR  (Markdown for a local RAG index)
# =============================================================================
def _md_escape_cell(text):
    """Make a cell safe inside a Markdown table: pipes escaped, newlines turned
    into <br> so a multi-line cell cannot break the row."""
    return (text or "").replace("|", "\\|").replace("\n", "<br>").strip()


def _render_markdown(prs, source_path):
    """
    Render a deck as Markdown for the RAG index: one '## Slide N' section per
    slide carrying its title, its bullets (indented by outline level), its
    tables and its speaker notes.

    Speaker notes are included deliberately. In a deck built to the 10/20/30
    rule the SLIDE carries the headline and the NOTES carry the argument, so a
    knowledge base that indexed only the slides would capture the labels and
    lose the content.
    """
    base = os.path.splitext(os.path.basename(source_path))[0]
    stamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
    slides = list(prs.slides)
    out = [
        "# {}".format(base),
        "",
        "- Source: {}".format(source_path),
        "- Slides: {}".format(len(slides)),
        "- Fetched: {}".format(stamp),
        "",
        "---",
        "",
    ]
    for i, slide in enumerate(slides):
        title = _slide_title(slide)
        out.append("## Slide {}{}".format(i + 1, " - " + title if title else ""))
        out.append("")
        try:
            out.append("*Layout: {}*".format(slide.slide_layout.name))
            out.append("")
        except (AttributeError, ValueError):
            pass

        for shape in slide.shapes:
            # The title is already the heading; don't repeat it as a bullet.
            if _ph_type_name(shape) in _TITLE_PH_TYPES:
                continue
            rows = _table_rows(shape)
            if rows:
                header = rows[0] if rows else []
                if header:
                    out.append("| " + " | ".join(_md_escape_cell(c) for c in header) + " |")
                    out.append("|" + "|".join(["---"] * len(header)) + "|")
                    for row in rows[1:]:
                        out.append("| " + " | ".join(_md_escape_cell(c) for c in row) + " |")
                    out.append("")
                continue
            if not _shape_text(shape).strip():
                continue
            try:
                paragraphs = shape.text_frame.paragraphs
            except (AttributeError, ValueError):
                continue
            for para in paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                out.append("{}- {}".format("  " * min(para.level, 8), text))
            out.append("")

        notes = _notes_text(slide)
        if notes:
            out.append("**Speaker notes:**")
            out.append("")
            for line in notes.split("\n"):
                out.append("> {}".format(line.strip()) if line.strip() else ">")
            out.append("")
    return "\n".join(out).rstrip() + "\n"


def _save_to_kb(prs, source_path):
    """
    Mirror the open deck to '<KB_DIR>/PowerPoint - <name>.md', overwriting any
    existing file, for a local RAG knowledge base. Returns the path written;
    raises OSError on failure.

    This is a side effect of the real operation; callers must never let a
    failure here break an open or a save, so failures are logged and swallowed
    by _mirror_to_kb.
    """
    base = os.path.splitext(os.path.basename(source_path))[0]
    content = _render_markdown(prs, source_path)
    os.makedirs(KB_DIR, exist_ok=True)
    filename = "PowerPoint - " + safe_filename(base) + ".md"
    path = os.path.join(KB_DIR, filename)
    # newline="\n" keeps endings consistent and avoids CRLF doubling on Windows.
    with open(path, "w", encoding="utf-8", newline="\n") as fh:
        fh.write(content)
    return path


def _mirror_to_kb(prs, source_path, result):
    """
    Mirror a deck to the knowledge-base folder if one is configured, and record
    the outcome on the tool's result dict. No-op when --kb-dir is unset.

    Called on open (what you READ enters the knowledge base) and on create and
    save (so does what you WRITE). Mirroring must never break the real
    operation, so failures are logged and reported, not raised.
    """
    if not KB_DIR:
        return
    try:
        kb_path = _save_to_kb(prs, source_path)
        log("saved to knowledge base: {}".format(kb_path))
        result["knowledge_base"] = kb_path
    except OSError as e:
        log("knowledge-base save failed: {}".format(e))
        result["knowledge_base_error"] = str(e)


# =============================================================================
# THE 10/20/30 AUDIT
# =============================================================================
# This is where the MCP server and the `kawasaki` skill meet: the skill explains
# the rule and shapes the content, and this measures whether the file on disk
# actually obeys it. Guidance alone drifts - a deck "designed for 30pt" that
# quietly inherits 24pt sub-bullets looks fine in the outline and fails in the
# room. Measuring the saved artefact is the only check that cannot be talked out
# of a finding.
def _estimate_minutes(prs):
    """
    Estimate speaking time from the SPEAKER NOTES word count, plus a fixed
    allowance per slide. Returns a dict carrying the estimate and every input
    that produced it, so the number can be argued with rather than trusted
    blindly.

    Notes rather than slide text, because slide text is the headline you show
    and the notes are the words you say. A deck with no notes cannot be timed;
    that is reported as such, not silently scored as zero.
    """
    slides = list(prs.slides)
    words = 0
    slides_with_notes = 0
    for slide in slides:
        notes = _notes_text(slide)
        if notes:
            slides_with_notes += 1
            words += len(notes.split())
    speaking = words / float(SPEAKING_WORDS_PER_MINUTE) if words else 0.0
    overhead = (len(slides) * SPEAKING_SECONDS_PER_SLIDE) / 60.0
    total = speaking + overhead
    return {
        "estimated_minutes": round(total, 1),
        "note_words": words,
        "slides_with_notes": slides_with_notes,
        "slides_without_notes": len(slides) - slides_with_notes,
        "words_per_minute": SPEAKING_WORDS_PER_MINUTE,
        "seconds_per_slide": SPEAKING_SECONDS_PER_SLIDE,
        "basis": (
            "speaker-note word count / {} wpm, plus {}s per slide".format(
                SPEAKING_WORDS_PER_MINUTE, SPEAKING_SECONDS_PER_SLIDE)
        ),
        "reliable": slides_with_notes > 0,
    }


def _audit_fonts(prs):
    """
    Every run of text whose EFFECTIVE size is below KAWASAKI_MIN_FONT_PT, with
    the slide, shape, resolved size and the inheritance step that set it.

    Returns (findings, smallest_pt, checked_runs, unresolved). `unresolved`
    collects TABLE text with no explicit size on its own run, paragraph or
    shape: such text is sized by the table style in tableStyles.xml, keyed by
    GUID, which this server does not read. Those are reported as unmeasured
    rather than guessed at in either direction - an honest "I could not measure
    this" beats both a false clean bill of health and a wave of false positives
    that buries the real findings.
    """
    findings = []
    unresolved = []
    smallest = None
    checked = 0
    for i, slide in enumerate(prs.slides):
        for shape, para, run, level, in_table in _iter_text_runs(slide):
            text = (run.text or "").strip()
            if not text:
                continue
            ph_type = _ph_type_name(shape)
            if ph_type in FONT_AUDIT_EXEMPT_PLACEHOLDERS:
                continue  # slide numbers/footers/dates are meant to be small
            pt, source = _effective_font_pt(prs, slide, shape, para, run, level)
            if in_table and source not in ("run", "paragraph", "shape"):
                # Nothing on the run, paragraph or shape itself sized this cell,
                # so the TABLE STYLE decides - and that lives in tableStyles.xml
                # keyed by GUID, which this server does not read. The chain
                # would happily hand back the master's otherStyle size here, but
                # that value does not govern table text: reporting it would be
                # false precision, and four "findings" per table would bury the
                # real ones. Report it as unmeasured instead.
                # The chain's answer is still worth showing as an ESTIMATE -
                # in the stock theme it does match what renders - but it is
                # labelled as one, because the step that produced it
                # (the master's otherStyle) is not what actually governs table
                # text. A number the user can sanity-check beats both a bare
                # "unknown" and a false assertion.
                unresolved.append({
                    "slide": i + 1,
                    "shape": shape.name,
                    "text": text[:80],
                    "estimated_pt": round(pt, 1),
                    "below_minimum": pt < KAWASAKI_MIN_FONT_PT,
                    "reason": "table text is sized by the deck's table style "
                              "(tableStyles.xml), which this server does not "
                              "read - estimated_pt is the surrounding text's "
                              "size, not a measurement of this cell",
                })
                continue
            checked += 1
            if smallest is None or pt < smallest:
                smallest = pt
            if pt < KAWASAKI_MIN_FONT_PT:
                findings.append({
                    "slide": i + 1,
                    "shape": shape.name,
                    "placeholder_type": ph_type,
                    "outline_level": level + 1,
                    "font_pt": round(pt, 1),
                    "source": source,
                    "text": text[:80],
                })
    findings.sort(key=lambda f: (f["font_pt"], f["slide"]))
    return findings, smallest, checked, unresolved


# =============================================================================
# TOOL HANDLERS  (each takes an `args` dict, returns a JSON-serialisable object)
# =============================================================================
def tool_list_presentations(args):
    """
    List the .pptx/.potx files available under the permitted roots, so the model
    can find a deck or a template by name instead of guessing a path. Each entry
    carries a 'location' ('docs' / 'templates') so a blank template is
    never mistaken for a real deck.
    """
    query = (args.get("query") or "").strip().lower()
    want_location = (args.get("location") or "").strip().lower()
    if want_location and want_location not in ("docs", "templates"):
        raise ToolError("location must be one of: docs, templates")
    items = []
    for p in _all_pptx_files():
        name = os.path.basename(p)
        if query and query not in name.lower():
            continue
        location = _root_label(p)
        if want_location and location != want_location:
            continue
        entry = {"name": name, "path": _display_path(p), "location": location}
        try:
            st = os.stat(p)
            entry["size_bytes"] = st.st_size
            entry["modified"] = datetime.datetime.fromtimestamp(
                st.st_mtime).isoformat(timespec="seconds")
        except OSError:
            pass
        items.append(entry)
    items.sort(key=lambda d: d["name"].lower())
    return {
        "count": len(items),
        "docs_dir": DOCS_DIR,
        "templates_dir": TEMPLATES_DIR,
        "presentations": items,
        "note": "Open any of these by passing its 'path' (or just its name) to "
                "powerpoint_open." + (
                    " Entries with location 'templates' are blank templates: "
                    "read-only, so start a new deck from one with "
                    "powerpoint_create(template=...) rather than editing it."
                    if TEMPLATES_DIR else ""),
    }


def _open_presentation(path):
    """Load a presentation, turning the engine's failures into clear messages."""
    try:
        return Presentation(path)
    except PackageNotFoundError:
        raise ToolError(
            "Not a valid .pptx file (corrupt, or not really Office Open XML): "
            "{}".format(path))
    except PermissionError:
        raise ToolError("Permission denied opening: {}".format(path))
    except Exception as e:
        raise ToolError("Could not open {}: {}".format(path, e))


def _register_session(path, prs):
    """Store an open presentation and return its session id."""
    if len(SESSIONS) >= MAX_SESSIONS:
        raise ToolError(
            "Too many open sessions ({}). Close some with "
            "powerpoint_close.".format(MAX_SESSIONS))
    sid = uuid.uuid4().hex[:12]
    SESSIONS[sid] = {"path": path, "prs": prs, "opened_at": _now_iso()}
    return sid


def tool_open(args):
    raw = _require(args, "path")
    path, fuzzy_matched = _resolve_existing_pptx(raw)
    prs = _open_presentation(path)
    sid = _register_session(path, prs)
    log("opened {} as session {}".format(path, sid))
    result = {
        "session_id": sid,
        "path": path,
        "slides": len(prs.slides),
        "layouts": len(prs.slide_layouts),
        "slide_size": "{} x {} in".format(
            round(prs.slide_width / 914400.0, 2),
            round(prs.slide_height / 914400.0, 2)),
    }
    if fuzzy_matched:
        # The opened file's name differs from what was asked for; surface it so
        # the caller can confirm it picked the intended deck.
        result["fuzzy_matched"] = True
        result["requested"] = raw
    if _root_label(path) == "templates":
        result["read_only"] = True
        result["note"] = ("This file is in the read-only templates folder. You "
                          "can read it and call powerpoint_list_layouts on it, "
                          "but every save is refused - build from it with "
                          "powerpoint_create(template=...) instead.")
    _mirror_to_kb(prs, path, result)
    return result


def tool_create(args):
    """
    Create a NEW .pptx in the presentations folder, open it as a session and
    return the session_id. The model then builds it up with
    powerpoint_add_slide and persists it with powerpoint_save.

    With no 'template', the stock Office template is used. With 'template' set
    to an existing .pptx/.potx (resolved within the permitted roots exactly as
    powerpoint_open resolves it), that file becomes the starting point - its
    masters, layouts, theme, fonts and colours are inherited - and the new deck
    is saved into the PRESENTATIONS folder, leaving the template untouched.

    There is no separate output folder: a deck Eva builds lands in the same
    library as the ones you gave it. To file it elsewhere afterwards, save-as
    with powerpoint_save(path=...) to any sub-folder of the presentation root.
    """
    filename = _require(args, "filename")
    if not isinstance(filename, str) or not filename.strip():
        raise ToolError("filename must be a non-empty string")

    # Confine to a bare filename: strip any directory components so the model
    # cannot traverse out of the presentations folder via '..' or an absolute
    # path.
    name = os.path.basename(filename.strip().replace("\\", "/"))
    if not name or name in (".", ".."):
        raise ToolError("filename must be a real file name, not a path")
    # A deck is always saved as .pptx - .potx is readable as a template but this
    # server does not author templates.
    stem, ext = os.path.splitext(name)
    if ext.lower() != ".pptx":
        name = stem + ".pptx"

    if not DOCS_DIR:
        raise ToolError("No presentations folder configured. Set --docs-dir "
                        "(or POWERPOINT_DOCS_DIR).")
    out_dir = os.path.realpath(os.path.expanduser(DOCS_DIR))
    try:
        os.makedirs(out_dir, exist_ok=True)
    except OSError as e:
        raise ToolError("Could not create folder {}: {}".format(out_dir, e))

    target = os.path.join(out_dir, name)
    # A presentations folder that contains (or IS) the read-only templates
    # folder would otherwise let a new deck be created among the blanks.
    _refuse_if_read_only(target, what="Creating a deck")

    if os.path.exists(target) and not bool(args.get("overwrite", False)):
        raise ToolError(
            "File already exists: {} (pass overwrite=true to replace "
            "it).".format(target))

    template = args.get("template")
    template_used = None
    if template is not None and str(template).strip():
        tmpl_path, _fuzzy = _resolve_existing_pptx(str(template))
        if os.path.realpath(tmpl_path) == os.path.realpath(target):
            raise ToolError("template and the new file resolve to the same path.")
        prs = _open_presentation(tmpl_path)
        template_used = _display_path(tmpl_path)
        # A template that already contains slides is a worked EXAMPLE, not a
        # blank. Keeping them would silently prepend someone else's content to
        # the new deck, so they go - the masters, layouts and theme (the part
        # that actually carries the styling) all stay.
        removed = _delete_all_slides(prs)
    else:
        prs = Presentation()
        removed = 0

    try:
        prs.save(target)
    except PermissionError:
        raise ToolError(
            "Permission denied creating {} (is it open in PowerPoint?).".format(target))
    except OSError as e:
        raise ToolError("Could not create {}: {}".format(target, e))

    sid = _register_session(target, prs)
    log("created {} as session {}{}".format(
        target, sid,
        " from template {}".format(template_used) if template_used else ""))
    rec = _recommended_layouts(prs)
    result = {
        "session_id": sid,
        "path": target,
        "created": True,
        "slides": len(prs.slides),
        "layouts": len(prs.slide_layouts),
        "recommended_layouts": rec,
        "next": "Call powerpoint_list_layouts to see this template's layouts "
                "and their font sizes, then powerpoint_add_slide for each "
                "slide, then powerpoint_save.",
    }
    if template_used:
        result["template"] = template_used
        if removed:
            result["template_example_slides_removed"] = removed
            result["note"] = (
                "The template contained {} example slide(s); they were removed "
                "so you start from a blank deck. Its masters, layouts, theme, "
                "fonts and colours are all inherited.".format(removed))
    _mirror_to_kb(prs, target, result)
    return result


def tool_close(args):
    sid = _require(args, "session_id")
    if sid not in SESSIONS:
        raise ToolError("Unknown session_id '{}'.".format(sid))
    SESSIONS.pop(sid)
    log("closed session {}".format(sid))
    return {"closed": sid, "note": "Any unsaved changes were discarded."}


def tool_list_sessions(args):
    out = []
    for sid, s in SESSIONS.items():
        out.append({
            "session_id": sid,
            "path": s["path"],
            "opened_at": s["opened_at"],
            "slides": len(s["prs"].slides),
        })
    return {"sessions": out, "count": len(out)}


def tool_list_layouts(args):
    """
    List the slide layouts this deck's template defines, each with its
    placeholders and the EFFECTIVE font size text in them will render at.

    This is the tool that makes "stick to the template styles" actionable: a
    branded template names its layouts whatever it likes, so the only way to use
    it correctly is to read what it actually offers rather than assume the stock
    Office deck's numbering.
    """
    session = _get_session(args)
    prs = session["prs"]
    query = (args.get("query") or "").strip().lower()
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        if query and query not in (layout.name or "").lower():
            continue
        if len(layouts) >= MAX_LAYOUTS_RETURNED:
            break
        placeholders = [_layout_placeholder_summary(prs, layout, ph)
                        for ph in layout.placeholders]
        layouts.append({
            "index": i,
            "name": layout.name,
            "role": _layout_role(layout),
            "placeholders": placeholders,
        })
    rec = _recommended_layouts(prs)
    # Which layouts can actually hold body text at the required size - the
    # question the 10/20/30 rule makes you ask before you write a word.
    compliant = []
    for entry in layouts:
        body = [p for p in entry["placeholders"]
                if p["type"] in _BODY_PH_TYPES and p["type"] != "SUBTITLE"]
        if body and all(p["level1_font_pt"] >= KAWASAKI_MIN_FONT_PT for p in body):
            compliant.append({"index": entry["index"], "name": entry["name"]})
    return {
        "count": len(layouts),
        "layouts": layouts,
        "recommended": rec,
        "min_font_pt": KAWASAKI_MIN_FONT_PT,
        "layouts_meeting_min_font_at_level_1": compliant,
        "note": "Pass a layout to powerpoint_add_slide by index, by name, or by "
                "role word ('title', 'section', 'bullets', 'two_content', "
                "'title_only', 'picture', 'blank'). 'level1_font_pt' is the "
                "size resolved through the whole inheritance chain, and "
                "'deeper_levels' shows what sub-bullets shrink to - that is "
                "usually where the {}-point rule breaks.".format(
                    KAWASAKI_MIN_FONT_PT),
    }


def _note_single_append(session):
    """
    Safety net for powerpoint_add_slide.

    It appends its slide to the END of the deck, so the deck's order is the
    order the CALLS arrive - and an MCP client may dispatch independent tool
    calls concurrently, in which case the arrival order is not the order the
    model wrote them. That is what shuffles a deck.

    We cannot see the intended order from inside a single call, but we CAN spot
    the shape of a parallel batch: several single appends landing on the same
    session within a fraction of a second, which no sequential turn-by-turn
    conversation produces. Returns a warning string (once per burst) or None.
    """
    now = time.monotonic()
    burst = session.get("append_burst")
    if burst is None or now - burst["last"] > APPEND_BURST_SECONDS:
        session["append_burst"] = {"last": now, "count": 1, "warned": False}
        return None
    burst["last"] = now
    burst["count"] += 1
    if burst["warned"] or burst["count"] < APPEND_BURST_MIN_CALLS:
        return None
    burst["warned"] = True
    return (
        "{} powerpoint_add_slide calls arrived on this session within {:.0f}ms "
        "of each other, which is what a batch of PARALLEL tool calls looks "
        "like. Each call appends to the end of the deck, so the slide order is "
        "whatever order the calls reached the server - not necessarily the order "
        "they were written in. Check the deck with powerpoint_get_content, fix "
        "the sequence with powerpoint_move_slide, and build decks with "
        "powerpoint_add_slides instead: one call, one ordered 'slides' list, "
        "order guaranteed.".format(burst["count"], APPEND_BURST_SECONDS * 1000)
    )


def _field(where, name):
    """
    Name a field for an error message: bare ("bullets") for a single-slide call,
    qualified ("slides[3].bullets") inside a batch, so a message always points
    at something the caller can actually find in what they sent.
    """
    return name if where is None else "{}.{}".format(where, name)


def _normalise_slide_spec(raw, where=None):
    """
    Validate one slide request and return a dict ready for _build_slide, or
    raise a ToolError naming the offending field. Everything that can be checked
    without touching the deck is checked HERE, so powerpoint_add_slides can
    validate a whole batch before adding a single slide.

    `where` is None for a single-slide call and "slides[N]" inside a batch.
    Note the layout is NOT resolved here: that needs the presentation, and
    _resolve_layout is called by the caller's own validation pass.
    """
    if not isinstance(raw, dict):
        raise ToolError(
            "{} must be an object describing one slide, e.g. "
            "{{\"layout\":\"bullets\",\"title\":\"...\",\"bullets\":[\"...\"]}} "
            "- got {}.".format(where or "a slide", type(raw).__name__))

    spec = {
        "layout": raw.get("layout"),
        "title": raw.get("title"),
        "subtitle": raw.get("subtitle"),
        "notes": raw.get("notes"),
        "placeholder": raw.get("placeholder"),
        "drop_empty_placeholders": bool(raw.get("drop_empty_placeholders", True)),
        "bullets": _normalise_bullets(
            raw.get("bullets"), argname=_field(where, "bullets")),
        "extra": [],
        "table": None,
    }

    # Extra placeholder fills - what a two-content layout's second column used
    # to need a follow-up powerpoint_set_placeholder call (and therefore a
    # slide_index) for. Keeping it in the same call removes both the extra round
    # trip and any chance of writing onto the wrong slide.
    extra = raw.get("placeholders")
    if extra is not None:
        if not isinstance(extra, list):
            raise ToolError(
                "{} must be a list of {{placeholder, text|bullets}} "
                "objects.".format(_field(where, "placeholders")))
        for j, item in enumerate(extra):
            at = "{}[{}]".format(_field(where, "placeholders"), j)
            if not isinstance(item, dict):
                raise ToolError("{} must be an object.".format(at))
            target = item.get("placeholder")
            if target is None:
                raise ToolError(
                    "{} needs a 'placeholder' (its idx, name or type "
                    "word).".format(at))
            if item.get("text") is not None:
                items = _normalise_bullets(
                    [str(item["text"])], argname="{}.text".format(at))
            elif item.get("bullets") is not None:
                items = _normalise_bullets(
                    item["bullets"], argname="{}.bullets".format(at))
            else:
                raise ToolError("{} needs either 'text' or 'bullets'.".format(at))
            spec["extra"].append({"placeholder": target, "items": items})

    table = raw.get("table")
    if table is not None:
        # Accept either {"rows": [[...]]} or the bare list of rows.
        rows_data = table.get("rows") if isinstance(table, dict) else table
        spec["table"] = _normalise_table_rows(
            rows_data, _field(where, "table"))

    return spec


def _build_slide(prs, spec, layout):
    """
    Append one slide on `layout` and fill it from a normalised spec. Returns
    (slide, result_entry). Shared by powerpoint_add_slide and
    powerpoint_add_slides so both behave identically.
    """
    slide = prs.slides.add_slide(layout)

    filled = []
    warnings = []

    title = spec.get("title")
    if title is not None and str(title).strip():
        try:
            ph = _placeholder_lookup(slide, "title")
        except ToolError:
            ph = None
        if ph is None:
            warnings.append(
                "This layout ('{}') has no title placeholder, so 'title' was "
                "ignored.".format(layout.name))
        else:
            _fill_text_frame(ph.text_frame, [{"text": str(title), "level": 0}])
            filled.append({"placeholder": ph.name, "type": _ph_type_name(ph)})

    subtitle = spec.get("subtitle")
    if subtitle is not None and str(subtitle).strip():
        ph = None
        for candidate in slide.placeholders:
            if _ph_type_name(candidate) == "SUBTITLE":
                ph = candidate
                break
        if ph is None:
            warnings.append(
                "This layout ('{}') has no subtitle placeholder, so 'subtitle' "
                "was ignored - a 'title' role layout usually has one.".format(
                    layout.name))
        else:
            _fill_text_frame(ph.text_frame, [{"text": str(subtitle), "level": 0}])
            filled.append({"placeholder": ph.name, "type": "SUBTITLE"})

    bullets = spec.get("bullets") or []
    if bullets:
        target = spec.get("placeholder")
        if target is None:
            # Default to the layout's first body/content placeholder.
            ph = None
            for candidate in slide.placeholders:
                if _ph_type_name(candidate) in ("BODY", "OBJECT"):
                    ph = candidate
                    break
            if ph is None:
                raise ToolError(
                    "Layout '{}' has no body/content placeholder to put bullets "
                    "in. Pick a layout with one (powerpoint_list_layouts shows "
                    "each layout's placeholders), or pass 'placeholder' "
                    "explicitly.".format(layout.name))
        else:
            ph = _placeholder_lookup(slide, target)
        _fill_text_frame(ph.text_frame, bullets)
        filled.append({"placeholder": ph.name, "type": _ph_type_name(ph),
                       "bullets": len(bullets)})

    # Extra placeholders (a two-content layout's second column, say). Filled
    # BEFORE the empty-placeholder sweep, so what they write is not swept away.
    for item in spec.get("extra") or []:
        ph = _placeholder_lookup(slide, item["placeholder"])
        written = _fill_text_frame(ph.text_frame, item["items"], replace=True)
        filled.append({"placeholder": ph.name, "type": _ph_type_name(ph),
                       "paragraphs": written})

    notes = spec.get("notes")
    if notes is not None and str(notes).strip():
        slide.notes_slide.notes_text_frame.text = str(notes)

    table_info = None
    if spec.get("table"):
        # Before the sweep as well: the table takes over an EMPTY body
        # placeholder's position, so that placeholder must still be there.
        table_info = _add_table_to_slide(prs, slide, spec["table"])

    removed = []
    if spec.get("drop_empty_placeholders", True):
        removed = _drop_empty_content_placeholders(slide)

    index = len(prs.slides) - 1
    entry = {
        "slide_index": index,
        "slide_number": index + 1,
        "layout": layout.name,
        "layout_index": list(prs.slide_layouts).index(layout),
        "filled": filled,
    }
    if table_info:
        entry["table"] = table_info
    if removed:
        entry["empty_placeholders_removed"] = removed
    if warnings:
        entry["warnings"] = warnings
    return slide, entry


def _kawasaki_slide_warning(prs):
    """The over-10-slides nudge, or None."""
    if len(prs.slides) <= KAWASAKI_MAX_SLIDES:
        return None
    return ("This deck now has {} slides, over the {}-slide guideline. Cut, or "
            "say why the extra slides earn their place.".format(
                len(prs.slides), KAWASAKI_MAX_SLIDES))


def tool_add_slide(args):
    """
    Append ONE slide built on one of the template's layouts, filling that
    layout's placeholders. Title, subtitle and bullets are all optional;
    whatever is left empty is removed by default so the deck carries no prompt
    text.

    Building a whole deck this way is a race - see _note_single_append and
    powerpoint_add_slides.
    """
    session = _get_session(args)
    prs = session["prs"]
    spec = _normalise_slide_spec(args)
    layout = _resolve_layout(prs, spec["layout"])
    _, entry = _build_slide(prs, spec, layout)

    result = dict(entry)
    result["slides"] = len(prs.slides)
    if result.get("table"):
        result["table"]["note"] = _TABLE_NOTE
    kawasaki = _kawasaki_slide_warning(prs)
    if kawasaki:
        result["kawasaki_warning"] = kawasaki
    order_warning = _note_single_append(session)
    if order_warning:
        result["order_warning"] = order_warning
    return result


MAX_BATCH_SLIDES = 200


def tool_add_slides(args):
    """
    Append an ORDERED list of slides in ONE call.

    This is the tool to build a deck with. powerpoint_add_slide appends its
    slide to the end of the deck, so the deck's order is the order the CALLS
    arrive; when a client dispatches a batch of them concurrently, the arrival
    order is not the order they were written in and the deck comes out shuffled.
    One call carrying the whole sequence cannot be reordered.
    """
    session = _get_session(args)
    prs = session["prs"]
    raw_slides = args.get("slides")
    if not isinstance(raw_slides, list) or not raw_slides:
        raise ToolError(
            "'slides' must be a non-empty list, in the order the slides should "
            "appear in the deck - e.g. [{\"layout\":\"title\",\"title\":\"...\"},"
            "{\"layout\":\"bullets\",\"title\":\"...\",\"bullets\":[\"...\"],"
            "\"notes\":\"...\"}].")
    if len(raw_slides) > MAX_BATCH_SLIDES:
        raise ToolError(
            "{} slides is more than this tool accepts in one call ({}). Split "
            "into consecutive calls - each one appends after the last, so the "
            "order still holds.".format(len(raw_slides), MAX_BATCH_SLIDES))

    # Validate EVERYTHING - including every layout - before adding any slide, so
    # a bad entry leaves the deck untouched instead of half-built.
    plan = []
    for i, raw in enumerate(raw_slides):
        where = "slides[{}]".format(i)
        spec = _normalise_slide_spec(raw, where)
        try:
            layout = _resolve_layout(prs, spec["layout"])
        except ToolError as e:
            raise ToolError("{}: {}".format(where, e))
        plan.append((spec, layout))

    added = []
    slides_before = len(prs.slides)
    for i, (spec, layout) in enumerate(plan):
        try:
            _, entry = _build_slide(prs, spec, layout)
        except ToolError as e:
            # A layout-vs-content mismatch (bullets onto a layout with no body
            # placeholder) can only be caught once the slide exists, so the
            # batch is unwound to leave the deck exactly as it was - including
            # the half-built slide add_slide() had already appended. Either the
            # whole batch lands or none of it does; a deck left half-written is
            # the state this tool exists to avoid.
            for index in range(len(prs.slides) - 1, slides_before - 1, -1):
                _remove_slide_at(prs, index)
            raise ToolError(
                "slides[{}] could not be built, so NOTHING was added (the deck "
                "still has its {} original slide(s)): {} Fix that entry and "
                "send the whole batch again.".format(i, slides_before, e))
        entry["index"] = i
        added.append(entry)

    # A batch is one arrival, so it can never be scrambled - reset the burst
    # tracker rather than letting a batch count towards it.
    session["append_burst"] = None
    log("add_slides: appended {} slides".format(len(added)))

    result = {"added": len(added), "slides_added": added,
              "slides": len(prs.slides)}
    if any(e.get("table") for e in added):
        result["table_note"] = _TABLE_NOTE
    kawasaki = _kawasaki_slide_warning(prs)
    if kawasaki:
        result["kawasaki_warning"] = kawasaki
    return result


def tool_set_placeholder(args):
    """Set the text of one placeholder on an existing slide - the tool for
    filling in a template's own title slide, or fixing a line."""
    session = _get_session(args)
    prs = session["prs"]
    slide = _get_slide(prs, _require(args, "slide_index"))
    ph = _placeholder_lookup(slide, _require(args, "placeholder"))
    if "text" in args and args["text"] is not None:
        items = _normalise_bullets([str(args["text"])], argname="text")
    else:
        items = _normalise_bullets(_require(args, "bullets"))
    written = _fill_text_frame(ph.text_frame, items, replace=True)
    return {
        "slide_index": _require(args, "slide_index"),
        "placeholder": ph.name,
        "type": _ph_type_name(ph),
        "paragraphs": written,
    }


def tool_add_bullets(args):
    """Append bullets to a placeholder that already has content, without
    clearing what is there."""
    session = _get_session(args)
    prs = session["prs"]
    slide = _get_slide(prs, _require(args, "slide_index"))
    spec = args.get("placeholder")
    if spec is None:
        ph = None
        for candidate in slide.placeholders:
            if _ph_type_name(candidate) in ("BODY", "OBJECT"):
                ph = candidate
                break
        if ph is None:
            raise ToolError(
                "This slide has no body/content placeholder. Pass "
                "'placeholder' explicitly, or use powerpoint_get_content to see "
                "what the slide has.")
    else:
        ph = _placeholder_lookup(slide, spec)
    items = _normalise_bullets(_require(args, "bullets"))
    written = _fill_text_frame(ph.text_frame, items, replace=False)
    return {
        "slide_index": _require(args, "slide_index"),
        "placeholder": ph.name,
        "bullets_added": written,
    }


def tool_set_notes(args):
    """
    Set a slide's speaker notes.

    Notes are load-bearing under the 10/20/30 rule, not an afterthought: the
    slide holds the headline and the notes hold the sentences you actually say,
    which is what keeps text off the slide and lets it stay at 30 points. They
    are also what powerpoint_review times the talk from.
    """
    session = _get_session(args)
    prs = session["prs"]
    index = _require(args, "slide_index")
    slide = _get_slide(prs, index)
    text = _require(args, "notes")
    if not isinstance(text, str):
        raise ToolError("notes must be a string")
    slide.notes_slide.notes_text_frame.text = text
    words = len(text.split())
    return {
        "slide_index": index,
        "note_words": words,
        "adds_seconds": round(words / float(SPEAKING_WORDS_PER_MINUTE) * 60, 1),
    }


_TABLE_NOTE = ("Table text is styled by the deck's table style, so this server "
               "does not resolve its font size - powerpoint_review reports "
               "table text as unmeasured rather than guessing.")


def _normalise_table_rows(rows_data, argname="rows"):
    """
    Validate a table's rows and return them as a list of lists of strings.
    Raises a ToolError naming `argname`, and touches nothing.
    """
    if not isinstance(rows_data, list) or not rows_data:
        raise ToolError("{} must be a non-empty list of row lists, e.g. "
                        "[[\"Region\",\"Q3\"],[\"APAC\",\"1.2m\"]]".format(argname))
    table_rows = []
    for r, row in enumerate(rows_data):
        if not isinstance(row, list):
            raise ToolError(
                "{}[{}] must be a list of cell strings.".format(argname, r))
        table_rows.append(["" if c is None else str(c) for c in row])
    if max(len(r) for r in table_rows) == 0:
        raise ToolError("{} must contain at least one column.".format(argname))
    return table_rows


def _add_table_to_slide(prs, slide, table_rows):
    """
    Draw a validated table (from _normalise_table_rows) onto `slide`. Where the
    layout has an empty body/content placeholder the table takes over that
    placeholder's exact position and size (and the placeholder is removed), so
    the table lands where the TEMPLATE says content goes rather than at some
    hard-coded offset. Returns a summary dict.
    """
    n_cols = max(len(r) for r in table_rows)
    n_rows = len(table_rows)

    # Prefer the template's own content area.
    host = None
    for candidate in slide.placeholders:
        if _ph_type_name(candidate) in ("BODY", "OBJECT", "TABLE") and \
                _is_empty_text_shape(candidate):
            host = candidate
            break
    if host is not None:
        left, top, width, height = host.left, host.top, host.width, host.height
        host._element.getparent().remove(host._element)
        placement = "the layout's content placeholder"
    else:
        # Fall back to a centred block with a margin proportional to the slide,
        # leaving the top third for the title.
        margin = int(prs.slide_width * 0.08)
        left = margin
        width = prs.slide_width - (2 * margin)
        top = int(prs.slide_height * 0.30)
        height = int(prs.slide_height * 0.55)
        placement = "a centred block (this layout had no free content placeholder)"

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    for r, row in enumerate(table_rows):
        for c in range(n_cols):
            # Setting cell.text keeps the table style's own fonts and colours.
            table.cell(r, c).text = row[c] if c < len(row) else ""
    return {"rows": n_rows, "columns": n_cols, "placement": placement}


def tool_add_table(args):
    """
    Add a table to a slide. Where the layout has an empty body/content
    placeholder, the table takes over that placeholder's exact position and size
    (and the placeholder is removed), so the table lands where the TEMPLATE says
    content goes rather than at some hard-coded offset.
    """
    session = _get_session(args)
    prs = session["prs"]
    index = _require(args, "slide_index")
    slide = _get_slide(prs, index)
    table_rows = _normalise_table_rows(_require(args, "rows"))
    info = _add_table_to_slide(prs, slide, table_rows)
    result = {"slide_index": index}
    result.update(info)
    result["note"] = _TABLE_NOTE
    return result


# --- Slide removal / reordering ----------------------------------------------
# python-pptx has no public API for either: a deck's slide ORDER lives in the
# presentation part's <p:sldIdLst>, and each entry points at a slide part by
# relationship id. Removing a slide therefore means dropping the relationship
# AND the sldId entry; reordering means moving the entry. Both are stable,
# well-understood OOXML operations - they are just not wrapped for us.
def _sld_id_list(prs):
    """The <p:sldIdLst> element that defines slide order."""
    el = prs.slides._sldIdLst
    if el is None:
        raise ToolError("This presentation has no slide list element.")
    return el


def _remove_slide_at(prs, index):
    """Remove the slide at zero-based `index`, dropping its relationship too so
    the saved package carries no orphaned part."""
    sld_id_lst = _sld_id_list(prs)
    entries = list(sld_id_lst)
    entry = entries[index]
    rid = entry.get(qn("r:id"))
    sld_id_lst.remove(entry)
    if rid:
        try:
            prs.part.drop_rel(rid)
        except (KeyError, AttributeError):
            # An already-dropped relationship is not worth failing over: the
            # sldIdLst entry is gone, so the slide is out of the deck either way.
            log("could not drop relationship {} (already gone?)".format(rid))


def _delete_all_slides(prs):
    """Strip every slide, leaving masters, layouts and theme intact. Used when a
    template ships with worked example slides. Returns how many went."""
    count = len(list(_sld_id_list(prs)))
    for i in range(count - 1, -1, -1):
        _remove_slide_at(prs, i)
    return count


def tool_delete_slide(args):
    """Delete one slide by zero-based index. Indices of later slides shift down
    by one, so delete from the highest index downwards when removing several."""
    session = _get_session(args)
    prs = session["prs"]
    index = _require(args, "slide_index")
    _get_slide(prs, index)  # validates the range and the type
    index = int(index)
    title = _slide_title(list(prs.slides)[index])
    _remove_slide_at(prs, index)
    return {
        "deleted_index": index,
        "deleted_title": title,
        "slides": len(prs.slides),
        "note": "Indices of later slides have shifted down by one. Delete from "
                "the highest index downwards when removing several.",
    }


def tool_move_slide(args):
    """Move a slide to a new position. Both indices are zero-based, and
    to_index is where the slide ends up in the FINAL ordering."""
    session = _get_session(args)
    prs = session["prs"]
    from_index = _require(args, "from_index")
    to_index = _require(args, "to_index")
    _get_slide(prs, from_index, "from_index")
    from_index = int(from_index)
    try:
        to_index = int(to_index)
    except (TypeError, ValueError):
        raise ToolError("to_index must be a whole number.")
    count = len(prs.slides)
    if not 0 <= to_index < count:
        raise ToolError(
            "to_index {} is out of range - this deck has {} slides "
            "(0-{}).".format(to_index, count, count - 1))
    sld_id_lst = _sld_id_list(prs)
    entry = list(sld_id_lst)[from_index]
    sld_id_lst.remove(entry)
    sld_id_lst.insert(to_index, entry)
    return {
        "from_index": from_index,
        "to_index": to_index,
        "order": [_slide_title(s) or "(untitled)" for s in prs.slides],
    }


def tool_get_content(args):
    """
    Read the deck back: every slide's layout, placeholder text, tables and
    speaker notes. This is how you check what a template's own slides contain
    before filling them in, and how you review a deck you did not build.
    """
    session = _get_session(args)
    prs = session["prs"]
    include_notes = bool(args.get("include_notes", True))
    only = args.get("slide_index")
    slides_out = []
    for i, slide in enumerate(prs.slides):
        if only is not None and int(only) != i:
            continue
        shapes_out = []
        for shape in slide.shapes:
            entry = {
                "name": shape.name,
                "placeholder_type": _ph_type_name(shape),
                "placeholder_idx": _ph_idx(shape),
            }
            rows = _table_rows(shape)
            if rows is not None:
                entry["table"] = rows
            else:
                text = _shape_text(shape)
                if not text.strip():
                    continue
                entry["text"] = text
                try:
                    entry["paragraphs"] = [
                        {"text": p.text, "level": p.level}
                        for p in shape.text_frame.paragraphs if p.text.strip()
                    ]
                except (AttributeError, ValueError):
                    pass
            shapes_out.append(entry)
        slide_entry = {
            "slide_index": i,
            "slide_number": i + 1,
            "title": _slide_title(slide),
            "shapes": shapes_out,
        }
        try:
            slide_entry["layout"] = slide.slide_layout.name
        except (AttributeError, ValueError):
            pass
        if include_notes:
            notes = _notes_text(slide)
            if notes:
                slide_entry["notes"] = notes
                slide_entry["note_words"] = len(notes.split())
        slides_out.append(slide_entry)
    if only is not None and not slides_out:
        raise ToolError(
            "slide_index {} is out of range - this deck has {} slides.".format(
                only, len(prs.slides)))
    return {
        "path": session["path"],
        "slides": len(prs.slides),
        "content": slides_out,
    }


def tool_review(args):
    """
    Audit the deck against Guy Kawasaki's 10/20/30 rule and report each part
    pass/fail with the evidence behind it.

    Deliberately reports rather than fixes. Two of the three findings have no
    mechanical fix: "too many slides" and "too long to present" are content
    decisions, and the font one is usually solved by choosing a different layout
    (or saying less), not by overriding a size - which would break the very
    template adherence this server exists to protect.
    """
    session = _get_session(args)
    prs = session["prs"]
    slides = list(prs.slides)
    n_slides = len(slides)

    timing = _estimate_minutes(prs)
    findings, smallest, checked, unresolved = _audit_fonts(prs)

    slides_pass = n_slides <= KAWASAKI_MAX_SLIDES
    minutes_pass = timing["estimated_minutes"] <= KAWASAKI_MAX_MINUTES
    font_pass = not findings

    # Slides carrying more text than an audience can read while listening. Not
    # part of the rule itself, but it is the failure the rule is aimed at, so it
    # is worth surfacing alongside.
    wordy = []
    for i, slide in enumerate(slides):
        words = 0
        for shape in slide.shapes:
            if _ph_type_name(shape) in _TITLE_PH_TYPES:
                continue
            words += len(_shape_text(shape).split())
        if words > 40:
            wordy.append({"slide": i + 1, "body_words": words,
                          "title": _slide_title(slide)})

    result = {
        "path": session["path"],
        "rule": "Guy Kawasaki's 10/20/30: {} slides, {} minutes, {}-point "
                "minimum font".format(KAWASAKI_MAX_SLIDES, KAWASAKI_MAX_MINUTES,
                                      KAWASAKI_MIN_FONT_PT),
        "passes": slides_pass and minutes_pass and font_pass,
        "slides": {
            "count": n_slides,
            "limit": KAWASAKI_MAX_SLIDES,
            "pass": slides_pass,
        },
        "minutes": {
            "limit": KAWASAKI_MAX_MINUTES,
            "pass": minutes_pass,
        },
        "font": {
            "minimum_pt": KAWASAKI_MIN_FONT_PT,
            "pass": font_pass,
            "smallest_pt": round(smallest, 1) if smallest is not None else None,
            "runs_checked": checked,
            "violations": len(findings),
            "findings": findings[:MAX_FONT_FINDINGS],
        },
    }
    result["minutes"].update(timing)

    if len(findings) > MAX_FONT_FINDINGS:
        result["font"]["truncated"] = (
            "Showing the {} smallest of {} findings.".format(
                MAX_FONT_FINDINGS, len(findings)))
    if unresolved:
        small_estimates = [u for u in unresolved if u["below_minimum"]]
        result["font"]["unmeasured"] = unresolved[:MAX_FONT_FINDINGS]
        result["font"]["unmeasured_note"] = (
            "Table text is sized by the deck's table style, which this server "
            "does not read, so these are NOT counted as violations. "
            + ("{} of them carry an estimate below {}pt - open the deck and "
               "check those tables by eye.".format(
                   len(small_estimates), KAWASAKI_MIN_FONT_PT)
               if small_estimates else
               "Their estimates all clear {}pt.".format(KAWASAKI_MIN_FONT_PT)))
        if small_estimates:
            todo_table = (
                "{} table cell(s) estimate below {}pt and could not be "
                "measured - check them by eye, and prefer fewer columns to a "
                "smaller font.".format(len(small_estimates), KAWASAKI_MIN_FONT_PT))
        else:
            todo_table = None
    else:
        todo_table = None
    if not timing["reliable"]:
        result["minutes"]["warning"] = (
            "No speaker notes, so the only input was {}s per slide. Add notes "
            "with powerpoint_set_notes to get a real estimate - the notes are "
            "the words you actually say.".format(SPEAKING_SECONDS_PER_SLIDE))
    if wordy:
        result["wordy_slides"] = wordy
        result["wordy_slides_note"] = (
            "These slides carry more than 40 words of body text. The 30-point "
            "rule exists to force this text into the speaker notes; a slide "
            "that still passes on size but reads like a document defeats it.")

    # A single actionable summary, so the caller does not have to reason over
    # three sub-objects to know what to do next.
    todo = []
    if not slides_pass:
        todo.append("Cut {} slide(s) to reach {}.".format(
            n_slides - KAWASAKI_MAX_SLIDES, KAWASAKI_MAX_SLIDES))
    if not minutes_pass:
        todo.append(
            "Trim roughly {} minute(s) of speaker notes to reach {}.".format(
                round(timing["estimated_minutes"] - KAWASAKI_MAX_MINUTES, 1),
                KAWASAKI_MAX_MINUTES))
    if not font_pass:
        smallest_finding = findings[0]
        todo.append(
            "{} run(s) below {}pt, smallest {}pt on slide {} (set by the "
            "'{}' step). Usually the fix is fewer words or a different layout, "
            "not a font override.".format(
                len(findings), KAWASAKI_MIN_FONT_PT, smallest_finding["font_pt"],
                smallest_finding["slide"], smallest_finding["source"]))
    if todo_table:
        todo.append(todo_table)
    result["todo"] = todo or ["Nothing to fix - the deck meets 10/20/30."]
    return result


def tool_save(args):
    """
    Save the deck. With no 'path' it is saved in place; with one, saved-as to
    that path (which must be inside the permitted folders and outside the
    read-only templates folder).
    """
    session = _get_session(args)
    prs = session["prs"]
    raw = args.get("path")
    if raw is None or not str(raw).strip():
        target = session["path"]
    else:
        target = _resolve_path(str(raw))
        stem, ext = os.path.splitext(target)
        if ext.lower() != ".pptx":
            # A deck is always .pptx; silently writing a .potx that PowerPoint
            # would then treat as a template is worse than correcting it.
            target = stem + ".pptx"
    _refuse_if_read_only(target)

    parent = os.path.dirname(target)
    if parent and not os.path.isdir(parent):
        try:
            os.makedirs(parent, exist_ok=True)
        except OSError as e:
            raise ToolError("Could not create folder {}: {}".format(parent, e))

    try:
        prs.save(target)
    except PermissionError:
        raise ToolError(
            "Permission denied saving {} (is it open in PowerPoint?).".format(target))
    except OSError as e:
        raise ToolError("Could not save {}: {}".format(target, e))

    session["path"] = target
    log("saved session to {}".format(target))
    result = {
        "saved": target,
        "slides": len(prs.slides),
    }
    _mirror_to_kb(prs, target, result)
    # Saving is the natural moment to say whether the deck is presentable, so
    # the 10/20/30 headline rides along rather than needing a second call.
    try:
        timing = _estimate_minutes(prs)
        findings, _smallest, _checked, _unresolved = _audit_fonts(prs)
        result["kawasaki"] = {
            "slides": len(prs.slides),
            "estimated_minutes": timing["estimated_minutes"],
            "runs_below_min_font": len(findings),
            "passes": (len(prs.slides) <= KAWASAKI_MAX_SLIDES
                       and timing["estimated_minutes"] <= KAWASAKI_MAX_MINUTES
                       and not findings),
            "detail": "Call powerpoint_review for the findings.",
        }
    except Exception as e:  # an audit must never fail a save
        log("post-save audit failed: {}".format(e))
    return result


# =============================================================================
# TOOL REGISTRY  (name -> (handler, JSON-Schema inputSchema, description))
# =============================================================================
# Shared guidance, so the wording cannot drift between tools.
_LAYOUT_HELP = (
    "Which layout to build the slide on: an index (0, 1, 2...), the layout's "
    "NAME as powerpoint_list_layouts reports it, or a role word - 'title', "
    "'section', 'bullets', 'two_content', 'title_only', 'picture', 'blank'. "
    "Role words are resolved against THIS template, so they work even when a "
    "branded deck names its layouts something else entirely. Omit it to use the "
    "template's own bullets/content layout."
)
_BULLETS_HELP = (
    "Bullet points, one per item: a list of strings, or a list of "
    "{\"text\": \"...\", \"level\": 0} objects where level is the 0-based "
    "outline depth (0 = top level, 1 = sub-bullet). Do NOT type '- ', '* ', "
    "a bullet glyph or '1. ' into the text - the layout supplies the marker, "
    "and a typed one is stripped. Every deeper level is SMALLER: check "
    "powerpoint_list_layouts' 'deeper_levels' before nesting, because level 2 "
    "is where the 30-point rule usually breaks."
)

TOOLS = [
    {
        "name": "powerpoint_list_presentations",
        "description": "List the .pptx/.potx files available under the presentation root (and the templates folder, if configured), with each file's path (relative to its root), 'location' ('docs' / 'templates'), size and last-modified time. Use this to find a deck or a TEMPLATE by name before calling powerpoint_open or powerpoint_create, instead of guessing paths - and with location='templates' to see the blank templates a new deck can be based on. Files in the templates folder are READ-ONLY: base a new deck on one with powerpoint_create(template=...) rather than opening and saving it. Optional 'query' filters the list by filename substring.",
        "handler": tool_list_presentations,
        "inputSchema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Optional case-insensitive filename substring filter, e.g. 'template'."},
                "location": {"type": "string", "enum": ["docs", "templates"], "description": "Optional: list only the files in one folder - 'docs' (the presentation root, which is also where new decks are created) or 'templates' (blank templates)."},
            },
        },
    },
    {
        "name": "powerpoint_create",
        "description": "Create a NEW .pptx in the server's presentations folder (--docs-dir) and open it as a session, returning the session_id every other tool needs. Pass 'template' to inherit an existing deck's masters, layouts, theme, fonts and colours - the single most important argument here, and the way to make output match a corporate deck. The template is a .pptx/.potx in the templates folder or the presentation root, resolved the same forgiving way as powerpoint_open (bare name, relative path, or a fuzzy near-miss); it is only READ and never modified, and any example slides in it are stripped so you start blank while keeping its styling. Without a template the stock Office design is used. Then call powerpoint_list_layouts to see what the template offers, ONE powerpoint_add_slides call carrying the whole ordered deck (a series of separate powerpoint_add_slide calls can arrive out of order and shuffle the slides), and powerpoint_save.",
        "handler": tool_create,
        "inputSchema": {
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "Name for the new file, e.g. 'series-a-pitch.pptx'. Any directory part is stripped; the file is created at the top of the configured presentations folder. The extension is forced to .pptx."},
                "template": {"type": "string", "description": "Optional: name or path of an existing .pptx/.potx whose design the new deck should inherit, e.g. 'Deck Template.pptx'. Looked up in the configured templates folder and the presentation root. The template is only read; the new deck is saved into the presentations folder."},
                "overwrite": {"type": "boolean", "default": False, "description": "Replace the file if it already exists (default false = error out)."},
            },
            "required": ["filename"],
        },
    },
    {
        "name": "powerpoint_open",
        "description": "Open an existing .pptx/.potx into an in-memory session and return a session_id used by all other tools. Just pass the file name (e.g. 'Kickoff.pptx') - it is resolved relative to the server's presentation root, so you do NOT need an absolute path; a relative sub-path and absolute/~ paths also work, and a bare name is found in subfolders. If no exact match is found it falls back to a FUZZY name match (so 'quarterly deck' can open 'Quarterly Deck 2024.pptx'); when a fuzzy match is used the result includes fuzzy_matched=true so you can confirm it opened the file you meant. Use this to READ a deck, or to inspect a template's layouts before drafting - but to BUILD from a template use powerpoint_create(template=...), because files in the templates folder can never be saved.",
        "handler": tool_open,
        "inputSchema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "The .pptx/.potx file name (relative to the presentation root), a relative sub-path, or an absolute/~ path."}
            },
            "required": ["path"],
        },
    },
    {
        "name": "powerpoint_list_layouts",
        "description": "List the slide LAYOUTS this deck's template defines - the tool that makes template adherence possible. For each layout it reports the index and name, a detected 'role' ('title' / 'section' / 'bullets' / 'two_content' / 'title_only' / 'picture' / 'blank'), and every placeholder with its idx, type, name and - the useful part - the EFFECTIVE font size text in it will render at, resolved through the full PowerPoint inheritance chain (run, paragraph, shape, layout, master, master txStyles, presentation default). Body placeholders also report 'deeper_levels', the sizes sub-bullets shrink to. The 'recommended' block names the best layout in THIS template for each job, and 'layouts_meeting_min_font_at_level_1' lists the ones whose body text clears the 30-point minimum. Call this straight after powerpoint_create, before adding any slide: a branded template names its layouts whatever it likes, so assuming the stock Office numbering is how a deck ends up off-template.",
        "handler": tool_list_layouts,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "query": {"type": "string", "description": "Optional case-insensitive substring filter on the layout name."},
            },
            "required": ["session_id"],
        },
    },
    {
        "name": "powerpoint_add_slides",
        "description": "Append MANY slides in ONE call, in the exact order given - THE TOOL TO BUILD A DECK WITH. Use it for anything past a single slide, and never fire a series of powerpoint_add_slide calls instead: each of those appends to the END of the deck, so the deck's order is the order the CALLS ARRIVE at the server, and independent tool calls issued together may be dispatched in parallel and arrive in any order - which comes out as a shuffled deck. One call carrying the whole sequence cannot be reordered. Each entry of 'slides' takes exactly what powerpoint_add_slide takes - layout, title, subtitle, bullets, placeholder, notes, drop_empty_placeholders - plus two things that used to need a second call against a slide_index: 'placeholders' (a list of {placeholder, text|bullets} for a two-content layout's other column) and 'table' ({rows: [[...]]}). Every entry, including every layout name, is validated BEFORE any slide is added, so a bad entry leaves the deck untouched and the error names its index. Returns each slide's real slide_index in order, so any later edit addresses the right slide.",
        "handler": tool_add_slides,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slides": {
                    "type": "array",
                    "minItems": 1,
                    "description": "The slides to append, IN DECK ORDER.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {"type": ["string", "integer"], "description": _LAYOUT_HELP},
                            "title": {"type": "string", "description": "Text for the layout's title placeholder. Make it the slide's ONE assertion, not a topic label."},
                            "subtitle": {"type": "string", "description": "Text for the layout's subtitle placeholder, if it has one (title layouts usually do)."},
                            "bullets": {"type": ["array", "string"], "items": {"type": ["string", "object"]}, "description": _BULLETS_HELP},
                            "placeholder": {"type": ["string", "integer"], "description": "Optional: which placeholder the bullets go in, by idx, name or type word. Defaults to the layout's first body/content placeholder."},
                            "placeholders": {"type": "array", "items": {"type": "object"}, "description": "Optional EXTRA placeholder fills, each {\"placeholder\": <idx|name|type word>, \"text\": \"...\"} or {\"placeholder\": ..., \"bullets\": [...]}. This is how a two-content layout's second column is filled without a separate powerpoint_set_placeholder call."},
                            "table": {"type": "object", "description": "Optional table for this slide: {\"rows\": [[\"Region\",\"Q3\"],[\"APAC\",\"1.2m\"]]}, first row treated as the header. It takes over an empty content placeholder's position, so it lands where the template says content goes."},
                            "notes": {"type": "string", "description": "Speaker notes for this slide - where the argument lives, so the slide can stay a headline at 30+ points. Also what powerpoint_review times the talk from."},
                            "drop_empty_placeholders": {"type": "boolean", "default": True, "description": "Remove TEXT placeholders left empty (default true). Picture/table/chart placeholders are never removed."},
                        },
                    },
                },
            },
            "required": ["session_id", "slides"],
        },
    },
    {
        "name": "powerpoint_add_slide",
        "description": "Append ONE slide built on one of the template's layouts and fill that layout's placeholders. Use this only for a single slide added on its own - to build a deck, or add any run of two or more slides, use powerpoint_add_slides, which keeps them in order (several add_slide calls issued together can be dispatched in parallel and arrive out of order, shuffling the deck). Text goes into PLACEHOLDERS, so the template's own fonts, sizes, colours and positions apply automatically; there is deliberately no way to set a font here, because that is what would break template adherence. Pass 'title', 'subtitle' and/or 'bullets' as the layout supports them, and 'notes' for the speaker notes (which is where the sentences you actually say belong). Unfilled TEXT placeholders are removed by default so the deck carries no 'Click to add text' prompts; picture/table/chart placeholders are always kept so a human can fill them in. Warns when the deck passes the 10-slide guideline.",
        "handler": tool_add_slide,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "layout": {"type": ["string", "integer"], "description": _LAYOUT_HELP},
                "title": {"type": "string", "description": "Text for the layout's title placeholder. Make it the slide's ONE assertion, not a topic label."},
                "subtitle": {"type": "string", "description": "Text for the layout's subtitle placeholder, if it has one (title layouts usually do)."},
                "bullets": {"type": ["array", "string"], "items": {"type": ["string", "object"]}, "description": _BULLETS_HELP},
                "placeholder": {"type": ["string", "integer"], "description": "Optional: which placeholder the bullets go in, by idx, name or type word ('body', 'content'). Defaults to the layout's first body/content placeholder - only needed for a two-content layout, where you name the second one explicitly."},
                "placeholders": {"type": "array", "items": {"type": "object"}, "description": "Optional EXTRA placeholder fills, each {\"placeholder\": <idx|name|type word>, \"text\": \"...\"} or {\"placeholder\": ..., \"bullets\": [...]} - a two-content layout's second column, filled in this same call rather than a follow-up powerpoint_set_placeholder."},
                "table": {"type": "object", "description": "Optional table for this slide: {\"rows\": [[\"Region\",\"Q3\"],[\"APAC\",\"1.2m\"]]}, first row treated as the header. It takes over an empty content placeholder's position, so it lands where the template says content goes."},
                "notes": {"type": "string", "description": "Speaker notes for this slide. Under the 10/20/30 rule this is where the argument lives, so the slide can stay a headline at 30+ points; it is also what powerpoint_review times the talk from."},
                "drop_empty_placeholders": {"type": "boolean", "default": True, "description": "Remove TEXT placeholders left empty, so no 'Click to add text' prompt remains (default true). Picture/table/chart placeholders are never removed. Set false to keep empty placeholders for a human to fill in PowerPoint."},
            },
            "required": ["session_id"],
        },
    },
    {
        "name": "powerpoint_set_placeholder",
        "description": "Set the text of ONE placeholder on an existing slide, replacing what is there. Use it to fill in a template's own title slide, correct a line, or write into the second placeholder of a two-content layout. Address the placeholder by its idx, its name, or a type word ('title', 'body', 'content', 'subtitle'). Pass 'text' for a single line, or 'bullets' for several. As everywhere in this server, no font is set - the placeholder and the outline level carry the styling. Take slide_index from powerpoint_get_content or from what powerpoint_add_slides returned, never from an assumption about the order slides were created in; and when the slide is one you are still building, fill its extra placeholders in the same powerpoint_add_slides entry ('placeholders') instead of following up here.",
        "handler": tool_set_placeholder,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "Zero-based slide index (powerpoint_get_content reports it)."},
                "placeholder": {"type": ["string", "integer"], "description": "Which placeholder: its idx (0, 1, ...), its exact name, or a type word - 'title', 'body', 'content', 'subtitle'."},
                "text": {"type": "string", "description": "A single line of text. Use 'bullets' instead for multiple paragraphs."},
                "bullets": {"type": ["array", "string"], "items": {"type": ["string", "object"]}, "description": _BULLETS_HELP},
            },
            "required": ["session_id", "slide_index", "placeholder"],
        },
    },
    {
        "name": "powerpoint_add_bullets",
        "description": "Append bullets to a placeholder that already has content, without clearing it. Use powerpoint_set_placeholder to replace instead. Defaults to the slide's first body/content placeholder. Take slide_index from powerpoint_get_content or from what powerpoint_add_slides returned - never from an assumption about creation order.",
        "handler": tool_add_bullets,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "Zero-based slide index."},
                "bullets": {"type": ["array", "string"], "items": {"type": ["string", "object"]}, "description": _BULLETS_HELP},
                "placeholder": {"type": ["string", "integer"], "description": "Optional: which placeholder, by idx, name or type word. Defaults to the slide's first body/content placeholder."},
            },
            "required": ["session_id", "slide_index", "bullets"],
        },
    },
    {
        "name": "powerpoint_set_notes",
        "description": "Set a slide's speaker notes, replacing any existing ones. Notes are not an afterthought under the 10/20/30 rule: the slide holds the headline and the notes hold the sentences you actually say, which is what keeps text off the slide and lets it stay at 30 points. powerpoint_review estimates the talk's length from these, so a deck with no notes cannot be timed. Returns the word count and the seconds it adds to the estimate.",
        "handler": tool_set_notes,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "Zero-based slide index."},
                "notes": {"type": "string", "description": "The speaker notes for this slide - what you will SAY, in full sentences. Pass an empty string to clear them."},
            },
            "required": ["session_id", "slide_index", "notes"],
        },
    },
    {
        "name": "powerpoint_add_table",
        "description": "Add a table to a slide from a list of row lists (the first row is the header). For a slide you are still building, pass the table in its powerpoint_add_slides entry instead, so it needs no slide_index at all. Where the layout has an empty body/content placeholder the table takes over that placeholder's exact position and size, so it lands where the TEMPLATE says content goes rather than at a hard-coded offset. Table text is styled by the deck's table style; this server does not override it, and powerpoint_review reports table text as unmeasured rather than guessing its size.",
        "handler": tool_add_table,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "Zero-based slide index."},
                "rows": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}, "description": "Rows of cell text, first row treated as the header, e.g. [[\"Region\",\"Q3\"],[\"APAC\",\"1.2m\"]]. Short rows are padded with empty cells."},
            },
            "required": ["session_id", "slide_index", "rows"],
        },
    },
    {
        "name": "powerpoint_get_content",
        "description": "Read the deck back: every slide's index, title, layout name, placeholder text (with each paragraph's outline level), tables and speaker notes. Use it to review a deck you did not build, to see what a template's own slides contain before filling them in, or to find the slide_index for an edit. Pass 'slide_index' to read one slide.",
        "handler": tool_get_content,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "Optional: read only this zero-based slide index."},
                "include_notes": {"type": "boolean", "default": True, "description": "Include speaker notes (default true)."},
            },
            "required": ["session_id"],
        },
    },
    {
        "name": "powerpoint_delete_slide",
        "description": "Delete one slide by zero-based index - the tool for getting a long deck back to ten slides. Indices of later slides shift down by one, so when removing several, delete from the HIGHEST index downwards (or re-read powerpoint_get_content between deletes).",
        "handler": tool_delete_slide,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "slide_index": {"type": "integer", "description": "Zero-based index of the slide to delete."},
            },
            "required": ["session_id", "slide_index"],
        },
    },
    {
        "name": "powerpoint_move_slide",
        "description": "Move a slide to a new position. Both indices are zero-based and 'to_index' is the position in the FINAL ordering. Returns the resulting slide order by title, so you can confirm the sequence reads as an argument rather than a list of topics.",
        "handler": tool_move_slide,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "from_index": {"type": "integer", "description": "Zero-based index of the slide to move."},
                "to_index": {"type": "integer", "description": "Zero-based position it should end up at."},
            },
            "required": ["session_id", "from_index", "to_index"],
        },
    },
    {
        "name": "powerpoint_review",
        "description": "Audit the deck against Guy Kawasaki's 10/20/30 rule and report each part pass/fail with its evidence: slide count against 10; estimated speaking time against 20 minutes (from the speaker-note word count, with every input shown so the number can be argued with); and every run of text whose EFFECTIVE font size is below 30 points, with the slide, shape, resolved size and which inheritance step set it. That last part is the one a human cannot easily do: in a well-built deck almost no run carries an explicit size, so checking the obvious property finds nothing wrong with a deck set entirely in 12pt. Also flags slides carrying more than 40 words of body text. Call it before you save, and again after any trim. It reports rather than fixes, because the fixes are content decisions - say less, cut a slide, or choose a layout with bigger body text - not font overrides, which would break the template adherence this server exists to protect.",
        "handler": tool_review,
        "inputSchema": {
            "type": "object",
            "properties": {"session_id": {"type": "string"}},
            "required": ["session_id"],
        },
    },
    {
        "name": "powerpoint_save",
        "description": "Save the deck - in place with no 'path', or save-as to a path inside the permitted folders. Saves into the read-only templates folder are refused. The result carries a 10/20/30 headline (slide count, estimated minutes, runs below the minimum font) so you find out immediately whether the deck is presentable; call powerpoint_review for the detail.",
        "handler": tool_save,
        "inputSchema": {
            "type": "object",
            "properties": {
                "session_id": {"type": "string"},
                "path": {"type": "string", "description": "Optional save-as path (relative to the presentation root, or absolute). Omit to save in place. The extension is forced to .pptx."},
            },
            "required": ["session_id"],
        },
    },
    {
        "name": "powerpoint_close",
        "description": "Close a session and free its deck. Unsaved changes are discarded.",
        "handler": tool_close,
        "inputSchema": {
            "type": "object",
            "properties": {"session_id": {"type": "string"}},
            "required": ["session_id"],
        },
    },
    {
        "name": "powerpoint_list_sessions",
        "description": "List all currently open presentation sessions.",
        "handler": tool_list_sessions,
        "inputSchema": {"type": "object", "properties": {}},
    },
]

_TOOL_BY_NAME = {t["name"]: t for t in TOOLS}


# =============================================================================
# JSON-RPC / MCP DISPATCH
# =============================================================================
def _jsonrpc_result(msg_id, result):
    return {"jsonrpc": "2.0", "id": msg_id, "result": result}


def _jsonrpc_error(msg_id, code, message):
    return {"jsonrpc": "2.0", "id": msg_id, "error": {"code": code, "message": message}}


def handle_message(msg):
    """
    Process one decoded JSON-RPC message. Returns a response dict, or None for
    notifications (which must not be answered).
    """
    method = msg.get("method")
    msg_id = msg.get("id")
    is_notification = "id" not in msg
    params = msg.get("params") or {}

    if method == "initialize":
        client_proto = params.get("protocolVersion") or PROTOCOL_VERSION_FALLBACK
        return _jsonrpc_result(msg_id, {
            "protocolVersion": client_proto,
            "capabilities": {"tools": {}},
            "serverInfo": {"name": SERVER_NAME, "version": SERVER_VERSION},
        })

    if method in ("notifications/initialized", "initialized"):
        return None  # notification, no response

    if method == "ping":
        return _jsonrpc_result(msg_id, {})

    if method == "tools/list":
        listed = [
            {"name": t["name"], "description": t["description"],
             "inputSchema": t["inputSchema"]}
            for t in TOOLS
        ]
        return _jsonrpc_result(msg_id, {"tools": listed})

    if method == "tools/call":
        name = params.get("name")
        arguments = params.get("arguments") or {}
        tool = _TOOL_BY_NAME.get(name)
        if tool is None:
            return _jsonrpc_result(msg_id, {
                "content": [{"type": "text", "text": "Unknown tool: {}".format(name)}],
                "isError": True,
            })
        try:
            result = tool["handler"](arguments)
            text = json.dumps(result, indent=2, ensure_ascii=False)
            return _jsonrpc_result(msg_id, {
                "content": [{"type": "text", "text": text}],
                "isError": False,
            })
        except ToolError as e:
            return _jsonrpc_result(msg_id, {
                "content": [{"type": "text", "text": "Error: {}".format(e)}],
                "isError": True,
            })
        except Exception as e:  # unexpected - log full trace to stderr, return summary
            log("UNEXPECTED in tool {}:\n{}".format(name, traceback.format_exc()))
            return _jsonrpc_result(msg_id, {
                "content": [{"type": "text", "text": "Internal error: {}".format(e)}],
                "isError": True,
            })

    # Unknown method.
    if is_notification:
        return None
    return _jsonrpc_error(msg_id, -32601, "Method not found: {}".format(method))


def serve():
    """Main stdio loop: read newline-delimited JSON-RPC, dispatch, reply."""
    log("starting")
    log("interpreter: {}".format(sys.executable))
    log(" / ".join("{} {}".format(k, v) for k, v in sorted(_versions().items())))
    if DOCS_DIR:
        log("path sandbox DOCS_DIR = {}".format(DOCS_DIR))
    if TEMPLATES_DIR:
        log("templates folder (read-only) = {}".format(TEMPLATES_DIR))
    else:
        log("no templates folder configured (set --templates-dir to add one)")
    if KB_DIR:
        log("knowledge-base mirroring enabled -> {}".format(KB_DIR))
    else:
        log("knowledge-base mirroring disabled (set --kb-dir to enable)")
    log("10/20/30 thresholds: {} slides / {} minutes / {} pt".format(
        KAWASAKI_MAX_SLIDES, KAWASAKI_MAX_MINUTES, KAWASAKI_MIN_FONT_PT))

    for line in sys.stdin:
        line = line.strip()
        if not line:
            continue
        try:
            msg = json.loads(line)
        except json.JSONDecodeError:
            # Cannot know the id; emit a parse error with null id per spec.
            sys.stdout.write(json.dumps(_jsonrpc_error(None, -32700, "Parse error")) + "\n")
            sys.stdout.flush()
            continue

        try:
            response = handle_message(msg)
        except Exception as e:  # last-resort guard so the server never dies
            log("FATAL in handle_message:\n{}".format(traceback.format_exc()))
            response = _jsonrpc_error(msg.get("id"), -32603, "Internal error: {}".format(e))

        if response is not None:
            sys.stdout.write(json.dumps(response, ensure_ascii=False) + "\n")
            sys.stdout.flush()

    log("stdin closed, exiting")


# =============================================================================
# SELF-TEST  (--check): offline round-trip so a single transfer can be validated
# =============================================================================
def _check_template(path, body_sz=3600, layout_name="Brand Content"):
    """
    Build a throwaway 'corporate template': the stock deck with its master's
    body style bumped to an unusual size and one layout renamed, plus an example
    slide. Those two changes are what make template inheritance PROVABLE - if
    the deck built from it reports that size and that layout name, the styling
    genuinely came from the template rather than from python-pptx's defaults.
    """
    prs = Presentation()
    master = prs.slide_masters[0]
    tx_styles = master.element.find(_P_TXSTYLES)
    body_style = tx_styles.find(_P_BODYSTYLE)
    lvl1 = body_style.find(qn("a:lvl1pPr"))
    lvl1.find(_A_DEFRPR).set("sz", str(body_sz))
    # Rename a layout via its cSld element (python-pptx exposes no setter).
    c_sld = prs.slide_layouts[1]._element.find(qn("p:cSld"))
    c_sld.set("name", layout_name)
    # An example slide, so the "template example slides are stripped" path runs.
    example = prs.slides.add_slide(prs.slide_layouts[1])
    example.shapes.title.text = "EXAMPLE - delete me"
    prs.save(path)


def run_check():
    import tempfile
    import shutil
    global DOCS_DIR, TEMPLATES_DIR, KB_DIR

    for name, value in sorted(_versions().items()):
        print("[check] {:<12}: {}".format(name, value))
    print("[check] interpreter : {}".format(sys.executable))

    failures = []

    def expect(label, condition, detail=""):
        ok = bool(condition)
        print("[check] {:<40} {}{}".format(
            label + ":", "PASS" if ok else "FAIL",
            "" if ok else "  <- " + str(detail)))
        if not ok:
            failures.append(label)
        return ok

    root = tempfile.mkdtemp(prefix="powerpoint_check_")
    try:
        # One documents folder holds both the source library and what the
        # server creates; templates are the read-only second root beside it.
        DOCS_DIR = os.path.join(root, "documents")
        TEMPLATES_DIR = os.path.join(root, "templates")
        KB_DIR = os.path.join(root, "knowledge")
        for folder in (DOCS_DIR, TEMPLATES_DIR, KB_DIR):
            os.makedirs(folder, exist_ok=True)
        SESSIONS.clear()

        tmpl_path = os.path.join(TEMPLATES_DIR, "Check Template.pptx")
        _check_template(tmpl_path)

        # --- the templates folder is visible and labelled -------------------
        listing = tool_list_presentations({})
        locations = {p["name"]: p["location"] for p in listing["presentations"]}
        expect("template is listed as 'templates'",
               locations.get("Check Template.pptx") == "templates", locations)

        # --- create from the template ---------------------------------------
        created = tool_create({"filename": "check-deck",
                               "template": "Check Template.pptx"})
        sid = created["session_id"]
        expect("create from template resolves it",
               created.get("template") == "Check Template.pptx", created.get("template"))
        expect("template example slides stripped",
               created["slides"] == 0 and
               created.get("template_example_slides_removed") == 1, created)
        expect("filename gets the .pptx extension",
               created["path"].endswith("check-deck.pptx"), created["path"])

        # --- the template's styling really was inherited ---------------------
        layouts = tool_list_layouts({"session_id": sid})
        brand = [lay for lay in layouts["layouts"] if lay["name"] == "Brand Content"]
        expect("renamed template layout is visible", len(brand) == 1,
               [lay["name"] for lay in layouts["layouts"]])
        body_pt = None
        if brand:
            bodies = [p for p in brand[0]["placeholders"]
                      if p["type"] in ("BODY", "OBJECT")]
            body_pt = bodies[0]["level1_font_pt"] if bodies else None
        expect("template's 36pt body size is resolved", body_pt == 36.0, body_pt)
        expect("layout roles detected",
               layouts["recommended"].get("title") is not None, layouts["recommended"])

        # --- build a small deck ---------------------------------------------
        s0 = tool_add_slide({
            "session_id": sid, "layout": "title",
            "title": "Ten slides, twenty minutes, thirty points",
            "subtitle": "A self-test deck",
            "notes": "Welcome. " * 40,
        })
        expect("title slide uses a title-role layout",
               s0["layout_index"] == layouts["recommended"]["title"]["index"], s0)
        expect("title slide fills both its placeholders",
               len(s0["filled"]) == 2, s0["filled"])

        s1 = tool_add_slide({
            "session_id": sid, "layout": "Brand Content",
            "title": "The problem",
            "bullets": ["- Typed markers are stripped",
                        {"text": "A sub-bullet shrinks", "level": 1}],
            "notes": "Here is the problem. " * 30,
        })
        expect("bullets land in the body placeholder",
               any(f.get("bullets") == 2 for f in s1["filled"]), s1["filled"])

        # --- typed bullet markers are stripped, levels are kept --------------
        content = tool_get_content({"session_id": sid, "slide_index": 1})
        paras = []
        for shape in content["content"][0]["shapes"]:
            paras.extend(shape.get("paragraphs", []))
        expect("typed '- ' marker stripped",
               any(p["text"] == "Typed markers are stripped" for p in paras), paras)
        expect("sub-bullet keeps outline level 1",
               any(p["level"] == 1 for p in paras), paras)

        # --- an unfilled TEXT placeholder is removed ------------------------
        # On its own scratch deck, so the main deck's slide count stays put.
        scratch = tool_create({"filename": "scratch-deck",
                               "template": "Check Template.pptx"})
        ssid = scratch["session_id"]
        bare = tool_add_slide({"session_id": ssid, "layout": "Brand Content",
                               "title": "Title only, body left empty"})
        expect("unfilled body placeholder is removed",
               bare.get("empty_placeholders_removed"), bare)
        kept = tool_add_slide({"session_id": ssid, "layout": "Brand Content",
                               "title": "Left for a human to fill",
                               "drop_empty_placeholders": False})
        expect("drop_empty_placeholders=false keeps it",
               "empty_placeholders_removed" not in kept, kept)
        scratch_slides = list(SESSIONS[ssid]["prs"].slides)
        expect("dropped placeholder is really gone from the slide",
               len(list(scratch_slides[0].placeholders)) == 1,
               [ph.name for ph in scratch_slides[0].placeholders])
        expect("kept placeholder is really still there",
               len(list(scratch_slides[1].placeholders)) == 2,
               [ph.name for ph in scratch_slides[1].placeholders])
        tool_close({"session_id": ssid})

        # --- notes, tables --------------------------------------------------
        notes_res = tool_set_notes({"session_id": sid, "slide_index": 1,
                                    "notes": "Replaced notes. " * 25})
        expect("notes word count reported", notes_res["note_words"] == 50, notes_res)
        table_res = tool_add_table({
            "session_id": sid, "slide_index": 0,
            "rows": [["Region", "Q3"], ["APAC", "1.2m"]]})
        expect("table added", table_res["rows"] == 2 and table_res["columns"] == 2,
               table_res)

        # --- the 30-point audit finds the inherited sub-bullet ---------------
        review = tool_review({"session_id": sid})
        sub_findings = [f for f in review["font"]["findings"]
                        if f["text"].startswith("A sub-bullet")]
        expect("audit flags the inherited sub-bullet size",
               len(sub_findings) == 1, review["font"]["findings"])
        expect("audit names the inheritance step",
               sub_findings and sub_findings[0]["source"] == "master_txstyles",
               sub_findings)
        expect("audit resolves inherited sizes, not just explicit ones",
               review["font"]["runs_checked"] > 0, review["font"])
        unmeasured_text = [u["text"] for u in review["font"].get("unmeasured", [])]
        expect("table text is reported unmeasured, not guessed at",
               sorted(unmeasured_text) == ["1.2m", "APAC", "Q3", "Region"],
               unmeasured_text)
        expect("unmeasured table text carries a labelled estimate",
               all("estimated_pt" in u and "below_minimum" in u
                   for u in review["font"]["unmeasured"]),
               review["font"]["unmeasured"][:1])
        expect("table text does not appear as a font violation",
               not [f for f in review["font"]["findings"]
                    if f["text"] in ("Region", "Q3", "APAC", "1.2m")],
               review["font"]["findings"])
        expect("slide count passes at 2 slides", review["slides"]["pass"], review["slides"])
        expect("timing estimated from notes",
               review["minutes"]["reliable"] and review["minutes"]["note_words"] > 0,
               review["minutes"])
        expect("review reports an actionable todo",
               isinstance(review["todo"], list) and review["todo"], review["todo"])

        # --- save, mirror, reopen -------------------------------------------
        saved = tool_save({"session_id": sid})
        expect("saved into the presentations folder",
               os.path.isfile(saved["saved"]) and
               _root_label(saved["saved"]) == "docs", saved)
        expect("save carries a 10/20/30 headline", "kawasaki" in saved, saved)
        kb_file = os.path.join(KB_DIR, "PowerPoint - check-deck.md")
        expect("knowledge-base mirror written", os.path.isfile(kb_file), kb_file)
        if os.path.isfile(kb_file):
            kb_text = open(kb_file, encoding="utf-8").read()
            expect("mirror carries slide titles",
                   "The problem" in kb_text, kb_text[:200])
            expect("mirror carries speaker notes",
                   "Replaced notes." in kb_text, kb_text[:200])
            expect("mirror carries tables", "| Region | Q3 |" in kb_text, kb_text[:400])

        tool_close({"session_id": sid})
        reopened = tool_open({"path": "check-deck.pptx"})
        rid = reopened["session_id"]
        expect("saved deck reopens with its slides", reopened["slides"] == 2, reopened)
        rc = tool_get_content({"session_id": rid})
        titles = [s["title"] for s in rc["content"]]
        expect("titles round-trip",
               titles == ["Ten slides, twenty minutes, thirty points", "The problem"],
               titles)
        expect("notes round-trip",
               rc["content"][1].get("notes", "").startswith("Replaced notes."),
               rc["content"][1].get("notes", "")[:60])

        # --- reorder and delete ---------------------------------------------
        moved = tool_move_slide({"session_id": rid, "from_index": 1, "to_index": 0})
        expect("slide moved to the front",
               moved["order"][0] == "The problem", moved["order"])
        deleted = tool_delete_slide({"session_id": rid, "slide_index": 0})
        expect("slide deleted", deleted["slides"] == 1, deleted)
        tool_save({"session_id": rid})
        after = tool_open({"path": "check-deck.pptx"})
        expect("deletion survives a save/reopen", after["slides"] == 1, after)
        tool_close({"session_id": after["session_id"]})
        tool_close({"session_id": rid})

        # --- the sandbox actually refuses ------------------------------------
        outside = os.path.join(root, "escape.pptx")
        try:
            _resolve_path(outside)
            expect("path outside the roots is refused", False, outside)
        except ToolError:
            expect("path outside the roots is refused", True)
        try:
            _resolve_path("../escape.pptx")
            expect("relative traversal is refused", False, "../escape.pptx")
        except ToolError:
            expect("relative traversal is refused", True)

        # --- templates are read-only -----------------------------------------
        tmpl_session = tool_open({"path": "Check Template.pptx"})
        expect("opening a template flags it read-only",
               tmpl_session.get("read_only") is True, tmpl_session)
        try:
            tool_save({"session_id": tmpl_session["session_id"]})
            expect("saving over a template is refused", False, "save succeeded")
        except ToolError:
            expect("saving over a template is refused", True)
        try:
            tool_save({"session_id": tmpl_session["session_id"],
                       "path": os.path.join(TEMPLATES_DIR, "sneaky.pptx")})
            expect("save-as into the templates folder is refused", False,
                   "save-as succeeded")
        except ToolError:
            expect("save-as into the templates folder is refused", True)
        tool_close({"session_id": tmpl_session["session_id"]})

        # --- a clean deck passes the whole rule -------------------------------
        clean = tool_create({"filename": "clean-deck",
                             "template": "Check Template.pptx"})
        csid = clean["session_id"]
        tool_add_slide({"session_id": csid, "layout": "title",
                        "title": "A clean deck", "subtitle": "Passes 10/20/30",
                        "notes": "A short opening. " * 20})
        tool_add_slide({"session_id": csid, "layout": "Brand Content",
                        "title": "One assertion per slide",
                        "bullets": ["Top-level bullets only"],
                        "notes": "The argument lives here. " * 30})
        clean_review = tool_review({"session_id": csid})
        expect("a compliant deck passes 10/20/30",
               clean_review["passes"] is True, clean_review["todo"])
        tool_close({"session_id": csid})

        # --- over-long decks are caught ---------------------------------------
        big = tool_create({"filename": "big-deck"})
        bsid = big["session_id"]
        warned = False
        for i in range(KAWASAKI_MAX_SLIDES + 2):
            res = tool_add_slide({"session_id": bsid, "layout": "title_only",
                                  "title": "Slide {}".format(i + 1)})
            warned = warned or ("kawasaki_warning" in res)
        expect("11th slide triggers the 10-slide warning", warned)
        big_review = tool_review({"session_id": bsid})
        expect("over-long deck fails the slide check",
               big_review["slides"]["pass"] is False, big_review["slides"])
        expect("no-notes deck reports an unreliable estimate",
               big_review["minutes"].get("reliable") is False,
               big_review["minutes"])
        tool_close({"session_id": bsid})

        # --- powerpoint_add_slides: one call, guaranteed deck order ----------
        ordered = tool_create({"filename": "ordered-deck",
                               "template": "Check Template.pptx"})
        osid = ordered["session_id"]
        oprs = SESSIONS[osid]["prs"]
        batch = tool_add_slides({"session_id": osid, "slides": [
            {"layout": "title", "title": "One", "subtitle": "Sub",
             "notes": "Opening remarks."},
            {"layout": "bullets", "title": "Two",
             "bullets": ["alpha", {"text": "nested", "level": 1}],
             "notes": "The argument."},
            {"layout": "section", "title": "Three"},
            {"layout": "bullets", "title": "Four",
             "table": {"rows": [["Region", "Q3"], ["APAC", "1.2m"]]}},
            {"layout": "bullets", "title": "Five"},
        ]})
        expect("add_slides adds every slide",
               batch["added"] == 5 and batch["slides"] == 5, batch)
        # THE point of the tool: the slides land in the order they were given.
        titles = [_slide_title(s) for s in oprs.slides]
        expect("add_slides preserves deck order",
               titles == ["One", "Two", "Three", "Four", "Five"], titles)
        expect("add_slides reports each slide's real index",
               [e["slide_index"] for e in batch["slides_added"]] == [0, 1, 2, 3, 4],
               batch["slides_added"])
        expect("add_slides keeps outline levels",
               [p.level for p in oprs.slides[1].placeholders[1].text_frame.paragraphs
                if p.text.strip()] == [0, 1],
               [(p.text, p.level) for p in
                oprs.slides[1].placeholders[1].text_frame.paragraphs])
        expect("add_slides writes speaker notes",
               oprs.slides[0].notes_slide.notes_text_frame.text ==
               "Opening remarks.",
               oprs.slides[0].notes_slide.notes_text_frame.text)
        expect("a slide's table is built in the same call",
               any(sh.has_table for sh in oprs.slides[3].shapes) and
               batch["slides_added"][3]["table"]["columns"] == 2,
               batch["slides_added"][3].get("table"))
        expect("add_slides carries the table caveat once",
               "table_note" in batch, sorted(batch))

        # A second call appends AFTER the first, so order survives batching.
        tool_add_slides({"session_id": osid,
                         "slides": [{"layout": "bullets", "title": "Six"}]})
        expect("a second batch appends after the first",
               [_slide_title(s) for s in oprs.slides][-1] == "Six",
               [_slide_title(s) for s in oprs.slides])

        # Validation is all-or-nothing: a bad entry adds no slides at all.
        before = len(oprs.slides)
        for bad, needle in (
            ([{"layout": "bullets", "title": "ok"},
              {"layout": "No Such Layout"}], "slides[1]"),
            ([{"layout": "bullets", "bullets": [{"level": 2}]}], "no 'text' key"),
            ([{"layout": "bullets", "table": {"rows": []}}], "non-empty list"),
            ([{"layout": "bullets", "placeholders": [{"text": "x"}]}],
             "needs a 'placeholder'"),
            (["not an object"], "must be an object"),
        ):
            try:
                tool_add_slides({"session_id": osid, "slides": bad})
                expect("bad batch refused: {}".format(needle), False, bad)
            except ToolError as e:
                expect("bad batch refused ({})".format(needle),
                       needle in str(e), str(e))
        expect("a refused batch adds no slides",
               len(oprs.slides) == before, len(oprs.slides))

        # A layout-vs-content mismatch can only surface once the slide exists,
        # so the batch must UNWIND rather than leave the deck half-written.
        try:
            tool_add_slides({"session_id": osid, "slides": [
                {"layout": "bullets", "title": "would be added"},
                {"layout": "blank", "title": "no body", "bullets": ["boom"]},
                {"layout": "bullets", "title": "never reached"}]})
            expect("mid-batch failure unwinds the batch", False, "no error")
        except ToolError as e:
            expect("mid-batch failure says nothing was added",
                   "NOTHING was added" in str(e), str(e))
        expect("mid-batch failure leaves the deck untouched",
               len(oprs.slides) == before, len(oprs.slides))
        try:
            tool_add_slides({"session_id": osid, "slides": []})
            expect("empty slides list refused", False, "no error raised")
        except ToolError as e:
            expect("empty slides list refused", "non-empty list" in str(e), str(e))

        # An extra placeholder fill lands, and survives the empty-placeholder
        # sweep that would otherwise have removed the placeholder it wrote into.
        two = tool_add_slides({"session_id": osid, "slides": [{
            "layout": 3,          # stock 'Two Content'
            "title": "Columns",
            "bullets": ["left one"],
            "placeholders": [{"placeholder": 2, "bullets": ["right one"]}],
        }]})
        col_texts = [_shape_text(sh) for sh in oprs.slides[-1].shapes]
        expect("extra placeholder fill lands on the same slide",
               any("right one" in t for t in col_texts) and
               any("left one" in t for t in col_texts), col_texts)
        expect("extra placeholder is reported as filled",
               len(two["slides_added"][0]["filled"]) == 3,
               two["slides_added"][0]["filled"])

        # The burst safety net: a rapid run of single add_slide calls is
        # reported, once, from the third onwards - and a batch never trips it.
        SESSIONS[osid]["append_burst"] = None
        burst = [tool_add_slide({"session_id": osid, "layout": "bullets",
                                 "title": "B{}".format(i)}) for i in range(4)]
        expect("first two rapid add_slide calls are not flagged",
               "order_warning" not in burst[0] and "order_warning" not in burst[1],
               [sorted(b) for b in burst[:2]])
        expect("a rapid add_slide burst is flagged",
               "order_warning" in burst[2] and
               "powerpoint_add_slides" in burst[2]["order_warning"],
               burst[2].get("order_warning"))
        expect("the burst is flagged only once",
               "order_warning" not in burst[3], burst[3].get("order_warning"))
        expect("a batched call is never treated as a burst",
               "order_warning" not in tool_add_slides(
                   {"session_id": osid,
                    "slides": [{"layout": "bullets", "title": "Batched"}]}),
               "batch tripped the burst detector")

        # The order survives a save and reopen - it is in the file, not just in
        # the session's object graph.
        osaved = tool_save({"session_id": osid})["saved"]
        tool_close({"session_id": osid})
        rsid = tool_open({"path": osaved})["session_id"]
        reopened = [_slide_title(s) for s in SESSIONS[rsid]["prs"].slides]
        expect("deck order survives save and reopen",
               reopened[:6] == ["One", "Two", "Three", "Four", "Five", "Six"],
               reopened)
        tool_close({"session_id": rsid})

    finally:
        SESSIONS.clear()
        shutil.rmtree(root, ignore_errors=True)

    print("[check]")
    if failures:
        print("[check] CHECKS FAILED ({}): {}".format(len(failures), ", ".join(failures)))
        return 1
    print("[check] ALL CHECKS PASSED")
    return 0


DISABLE_KEYWORDS = frozenset(("off", "none", "no", "false", "disabled"))


def _is_disabled(value):
    """True if a folder setting is one of the DISABLE_KEYWORDS.

    An MCP client can only pass strings, and a BLANK string is what it
    substitutes for a setting the user left empty - which the suite treats as
    "not configured", falling back to the default. So a keyword is needed to say
    "definitely off": --kb-dir off (or POWERPOINT_KB_DIR=off) disables Markdown
    mirroring, --templates-dir off disables the templates root.
    """
    return bool(value) and value.strip().lower() in DISABLE_KEYWORDS


def main():
    parser = argparse.ArgumentParser(
        description="PowerPoint (.pptx) python-pptx MCP stdio server. "
                    "With no arguments it runs as an MCP server on stdin/stdout."
    )
    parser.add_argument(
        "--check", action="store_true",
        help="Run an offline create/build/save/reopen/audit self-test and exit."
    )
    parser.add_argument(
        "--version", action="version",
        version="{0} {1}".format(SERVER_NAME, __version__)
    )
    parser.add_argument(
        "--docs-dir", default=os.environ.get("POWERPOINT_DOCS_DIR"), metavar="DIR",
        help="REQUIRED path sandbox: the server refuses to open or save any "
             "file outside this directory tree, and refuses to start without "
             "one. Falls back to the POWERPOINT_DOCS_DIR environment variable, "
             "then the DOCS_DIR config value (default: "
             "C:\\Eva\\documents\\powerpoint). It is also where "
             "powerpoint_create writes NEW decks - there is no separate output "
             "folder. The model chooses open/save paths, so an unconfined "
             "server could read/write any .pptx this account can."
    )
    parser.add_argument(
        "--templates-dir", default=os.environ.get("POWERPOINT_TEMPLATES_DIR"),
        metavar="DIR",
        help="Folder of blank .pptx/.potx templates. READ-ONLY: its files can be "
             "listed, opened and passed as powerpoint_create's 'template', but "
             "nothing can be saved over them. Falls back to the "
             "POWERPOINT_TEMPLATES_DIR environment variable, then the "
             "TEMPLATES_DIR config value (default: "
             "C:\\Eva\\templates\\powerpoint); pass 'off' to run with no "
             "templates root at all."
    )
    parser.add_argument(
        "--kb-dir", default=os.environ.get("POWERPOINT_KB_DIR"), metavar="DIR",
        help="Every deck opened, created or saved is ALSO written as a Markdown "
             "file into this folder for a local RAG knowledge base (falls back "
             "to the POWERPOINT_KB_DIR environment variable, then the KB_DIR "
             "config value, default: C:\\Eva\\knowledge\\powerpoint); pass 'off' "
             "to disable mirroring. Files are named 'PowerPoint - <name>.md' and "
             "overwritten each time. The folder is created if it does not exist."
    )
    args = parser.parse_args()

    global DOCS_DIR, TEMPLATES_DIR, KB_DIR

    # Each optional folder has a real default (the Eva working tree), so "leave
    # it blank" cannot mean "turn this off": a blank value from an MCP client
    # means "not configured", and falls back to the default. Pass one of the
    # DISABLE_KEYWORDS instead to switch a folder off explicitly - the only way
    # to do it from a client that can only supply strings.
    # explicit_folders records which paths the user actually chose, so a typo in
    # one still fails loudly while a default pointing at a folder this endpoint
    # has not created yet degrades quietly.
    explicit_folders = set()
    for name, value in (("docs", args.docs_dir),
                        ("templates", args.templates_dir),
                        ("kb", args.kb_dir)):
        if not value:
            continue                      # unset -> keep the CONFIG default
        chosen = None if _is_disabled(value) else value
        explicit_folders.add(name)
        if name == "docs":
            DOCS_DIR = chosen
        elif name == "templates":
            TEMPLATES_DIR = chosen
        else:
            KB_DIR = chosen

    if args.check:
        sys.exit(run_check())

    # File access is confined to DOCS_DIR, so a root is mandatory.
    if not DOCS_DIR:
        log("FATAL: no presentation root configured. Pass --docs-dir, set the "
            "POWERPOINT_DOCS_DIR environment variable, or set the DOCS_DIR "
            "constant in this file. The server only opens/saves .pptx files "
            "inside that folder and will not start without one.")
        sys.exit(2)
    if not os.path.isdir(DOCS_DIR):
        log("FATAL: the configured presentation root does not exist or is not a "
            "directory: {}".format(DOCS_DIR))
        if "docs" not in explicit_folders:
            log("       That is the built-in default. Create the folder, or "
                "copy the repo's eva\\ folder to C:\\Eva to lay the whole "
                "working tree out at once, or point --docs-dir somewhere else.")
        sys.exit(2)

    # The templates folder is optional. A path the USER chose and got wrong
    # would silently mean "no templates", so that fails loudly; the built-in
    # default simply not existing yet does not.
    if TEMPLATES_DIR and not os.path.isdir(TEMPLATES_DIR):
        if "templates" in explicit_folders:
            log("FATAL: the configured templates folder does not exist or is "
                "not a directory: {}".format(TEMPLATES_DIR))
            sys.exit(2)
        log("WARNING: the default templates folder does not exist, so "
            "templates are disabled: {}".format(TEMPLATES_DIR))
        log("         Create it (copying the repo's eva\\ folder to C:\\Eva "
            "lays out the whole tree) to create decks from a template.")
        TEMPLATES_DIR = None

    if TEMPLATES_DIR:
        # Templates are read-only, so the folder must not be - or contain - the
        # presentation root, which is the one folder decks are saved into; that
        # would refuse every save.
        real_docs = os.path.realpath(os.path.expanduser(DOCS_DIR))
        if _read_only_root(real_docs) is not None:
            log("FATAL: the templates folder ({}) is the same as - or contains "
                "- the presentation root ({}). Templates are read-only, so "
                "every save would be refused. Point --templates-dir at a "
                "folder of its own.".format(TEMPLATES_DIR, DOCS_DIR))
            sys.exit(2)

    try:
        serve()
    except KeyboardInterrupt:
        log("interrupted, exiting")


if __name__ == "__main__":
    main()
